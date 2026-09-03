import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# ==============================================================================
# WORLD-CLASS ART DIRECTION: COLOR SYSTEM (Apple Keynote × Calm Tech)
# ==============================================================================
C_BG_HERO    = RGBColor(6, 10, 20)       # #060A14 Near-black deep cinematic navy
C_BG_DARK    = RGBColor(10, 15, 29)      # #0A0F1D Rich dark slate
C_SURFACE    = RGBColor(16, 24, 44)      # #10182C Surface container
C_SURFACE_HI = RGBColor(22, 33, 60)      # #16213C Elevated container
C_BORDER     = RGBColor(34, 48, 80)      # #223050 Minimal clean border line
C_BORDER_HI  = RGBColor(56, 78, 128)     # #384E80 Highlight border

C_CYAN       = RGBColor(6, 182, 212)     # #06B6D4 Electric clean cyan
C_TEAL       = RGBColor(20, 184, 166)    # #14B8A6 Soft calm teal
C_INDIGO     = RGBColor(99, 102, 241)    # #6366F1 Modern indigo
C_EMERALD    = RGBColor(16, 185, 129)    # #10B981 Vitality green
C_AMBER      = RGBColor(245, 158, 11)    # #F59E0B Warning/warmth
C_ROSE       = RGBColor(244, 63, 94)     # #F43F5E Coral/alert

C_TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF Crisp pure white
C_TEXT_LIGHT = RGBColor(226, 232, 240)   # #E2E8F0 Off-white body
C_TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 Silver secondary
C_TEXT_DIM   = RGBColor(100, 116, 139)   # #64748B Tertiary slate

FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"

def build_cinematic_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    logo_path = os.path.abspath("frontend/public/humsaathi-logo-v1.png")
    has_logo = os.path.exists(logo_path)

    shots = {
        "child": os.path.abspath("screenshots/child_portal.png"),
        "teen": os.path.abspath("screenshots/teen_portal.png"),
        "adult": os.path.abspath("screenshots/adult_portal.png"),
        "scenarios": os.path.abspath("screenshots/scenarios_portal.png"),
        "conversation": os.path.abspath("screenshots/conversation_portal.png"),
        "parent": os.path.abspath("screenshots/parent_portal.png"),
        "landing": os.path.abspath("screenshots/landing.png"),
        "login": os.path.abspath("screenshots/login_rendered.png"),
    }

    # Helper: Slide background with smooth transition
    def add_slide_base(slide, bg_color=C_BG_DARK, is_morph=False):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color
        bg.line.fill.background()

        if is_morph:
            xml = (
                '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                'xmlns:p16="http://schemas.microsoft.com/office/powerpoint/2015/09/main">'
                '<p16:morph option="byObject"/></p:transition>'
            )
        else:
            xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
        slide._element.append(parse_xml(xml))
        return bg

    def add_minimal_footer(slide, slide_num):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"HumSaathi AI  ·  Product Pitch Deck  ·  {slide_num:02d} / 15"
        p.font.name = FONT_BODY
        p.font.size = Pt(9)
        p.font.color.rgb = C_TEXT_DIM

    def add_headline(slide, tag, title, top=Inches(0.65)):
        # Minimal tag
        tb_tag = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.733), Inches(0.35))
        tf_tag = tb_tag.text_frame
        tf_tag.margin_left = tf_tag.margin_top = tf_tag.margin_right = tf_tag.margin_bottom = 0
        p_t = tf_tag.paragraphs[0]
        p_t.text = tag.upper()
        p_t.font.name = FONT_BODY
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = C_CYAN

        # Bold modern headline
        tb_h = slide.shapes.add_textbox(Inches(0.8), top + Inches(0.35), Inches(11.733), Inches(0.8))
        tf_h = tb_h.text_frame
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.name = FONT_HEAD
        p_h.font.size = Pt(32)
        p_h.font.bold = True
        p_h.font.color.rgb = C_TEXT_WHITE

    def add_browser_mockup(slide, img_path, left, top, width, height, title="humsaathi.ai"):
        # Browser header bar
        bar_h = Inches(0.32)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height + bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_SURFACE
        bar.line.color.rgb = C_BORDER
        bar.line.width = Pt(1)

        # Traffic light dots
        dot_colors = [RGBColor(239, 68, 68), RGBColor(245, 158, 11), RGBColor(16, 185, 129)]
        for i, dc in enumerate(dot_colors):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18 + i * 0.16), top + Inches(0.1), Inches(0.1), Inches(0.1))
            dot.fill.solid()
            dot.fill.fore_color.rgb = dc
            dot.line.fill.background()

        # URL text in browser bar
        url_box = slide.shapes.add_textbox(left + Inches(0.8), top + Inches(0.06), width - Inches(1.6), Inches(0.24))
        p = url_box.text_frame.paragraphs[0]
        p.text = f"https://{title}"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_BODY
        p.font.size = Pt(8.5)
        p.font.color.rgb = C_TEXT_DIM

        # Screenshot inside
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, left + Inches(0.02), top + bar_h, width=width - Inches(0.04), height=height)
            except Exception:
                pass

    # ==========================================================================
    # SLIDE 01: CINEMATIC OPENING (Full-Screen Product Launch Hero)
    # ==========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_slide_base(s1, C_BG_HERO)

    # Ambient subtle gradient glow rect
    glow = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.0), Inches(1.5), Inches(9.333), Inches(4.5))
    glow.fill.solid()
    glow.fill.fore_color.rgb = C_SURFACE
    glow.line.color.rgb = C_BORDER
    glow.line.width = Pt(1)

    # Official Logo centered
    if has_logo:
        try:
            s1.shapes.add_picture(logo_path, Inches(5.916), Inches(2.0), width=Inches(1.5))
        except Exception:
            pass

    # Huge cinematic typography
    h_box = s1.shapes.add_textbox(Inches(1.0), Inches(3.45), Inches(11.333), Inches(1.2))
    tf = h_box.text_frame
    p = tf.paragraphs[0]
    p.text = "HUMSAATHI AI"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEAD
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_WHITE

    tag_box = s1.shapes.add_textbox(Inches(1.0), Inches(4.55), Inches(11.333), Inches(0.5))
    p_tag = tag_box.text_frame.paragraphs[0]
    p_tag.text = "Learn. Communicate. Grow."
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.name = FONT_BODY
    p_tag.font.size = Pt(20)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_CYAN

    sub_box = s1.shapes.add_textbox(Inches(1.0), Inches(5.05), Inches(11.333), Inches(0.4))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = "The Adaptive AI Communication & Learning Platform"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = C_TEXT_MUTED

    # Minimal persona indicator line
    p_line = s1.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.333), Inches(0.4))
    p_l = p_line.text_frame.paragraphs[0]
    p_l.text = "CHILD  ·  TEEN  ·  ADULT"
    p_l.alignment = PP_ALIGN.CENTER
    p_l.font.name = FONT_BODY
    p_l.font.size = Pt(11)
    p_l.font.bold = True
    p_l.font.color.rgb = C_TEXT_DIM

    add_minimal_footer(s1, 1)

    s1.notes_slide.notes_text_frame.text = (
        "SLIDE 1 — CINEMATIC HERO (HumSaathi AI)\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Good morning, judges. Over 1 billion people worldwide navigate speech, language, and neurodiverse learning challenges. "
        "Yet modern education still treats learning as a one-way street.\n"
        "This is HumSaathi AI — 'Our Companion'. An adaptive platform designed to bridge the gap between knowing words "
        "and communicating with real-world confidence across the human lifespan: Child, Teen, and Adult.\"\n\n"
        "KEY JUDGE TAKEAWAY: Immediate emotional hook, serious AI product demeanor, clear lifespan scope."
    )

    # ==========================================================================
    # SLIDE 02: THE PROBLEM (Typography-Driven Visual Transformation)
    # ==========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_slide_base(s2, C_BG_DARK)

    # Massive Typography Statement
    big_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.733), Inches(2.2))
    b_tf = big_box.text_frame
    b_tf.word_wrap = True
    p1 = b_tf.paragraphs[0]
    p1.text = "LEARNING IS"
    p1.font.name = FONT_HEAD
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = C_TEXT_DIM

    p2 = b_tf.add_paragraph()
    p2.text = "NOT ONE-SIZE-FITS-ALL."
    p2.font.name = FONT_HEAD
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = C_ROSE

    # Horizontal divider line
    div = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.733), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = C_BORDER
    div.line.fill.background()

    # 4 Clean Columns of the Problem (No Cards, Pure Typography)
    problems = [
        ("01", "STATIC CONTENT", "Traditional apps force every learner through identical linear lessons, ignoring sensory sensitivities and processing pace."),
        ("02", "PRACTICE GAP", "Learners can memorize vocabulary in isolation, but have zero safe, judgment-free space to practice spoken dialogue."),
        ("03", "FEEDBACK VOID", "Standard educational software returns binary pass/fail checks, offering no coaching on tone, social nuance, or clarity."),
        ("04", "COMMUNICATION ANXIETY", "Facing unpredictable social and workplace interactions without simulation causes withdrawal and cognitive fatigue."),
    ]

    for i, (num, title, desc) in enumerate(problems):
        col_x = Inches(0.8) + i * Inches(2.98)
        c_box = s2.shapes.add_textbox(col_x, Inches(3.9), Inches(2.75), Inches(2.8))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True

        p_num = c_tf.paragraphs[0]
        p_num.text = num
        p_num.font.name = FONT_BODY
        p_num.font.size = Pt(13)
        p_num.font.bold = True
        p_num.font.color.rgb = C_CYAN

        p_t = c_tf.add_paragraph()
        p_t.text = title
        p_t.font.name = FONT_HEAD
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_WHITE
        p_t.space_before = Pt(4)

        p_d = c_tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(6)

    add_minimal_footer(s2, 2)

    s2.notes_slide.notes_text_frame.text = (
        "SLIDE 2 — THE PROBLEM (Learning Is Not One-Size-Fits-All)\n"
        "TIME: 30 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Why does communication learning break down for neurodiverse individuals? Four key failures:\n"
        "Static content that treats everyone identically.\n"
        "A gaping practice gap between knowing words and speaking them.\n"
        "A feedback void where right/wrong checkmarks teach zero conversational nuance.\n"
        "And communication anxiety that shuts down real-world independence.\n"
        "We set out to replace this broken model.\"\n\n"
        "KEY JUDGE TAKEAWAY: The problem is communication confidence and sensory adaptation, not simple literacy."
    )

    # ==========================================================================
    # SLIDE 03: THE BIG IDEA (AI Intelligence Visualization)
    # ==========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_slide_base(s3, C_BG_HERO)
    add_headline(s3, "The Central Intelligence", "Meet HumSaathi AI")

    # Center Hero Hub
    center_hub = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.666), Inches(2.2), Inches(4.0), Inches(4.2))
    center_hub.fill.solid()
    center_hub.fill.fore_color.rgb = C_SURFACE
    center_hub.line.color.rgb = C_CYAN
    center_hub.line.width = Pt(2)

    hub_box = s3.shapes.add_textbox(Inches(4.866), Inches(2.5), Inches(3.6), Inches(3.6))
    h_tf = hub_box.text_frame
    h_tf.word_wrap = True
    p = h_tf.paragraphs[0]
    p.text = "HUMSAATHI AI"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEAD
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    p_c = h_tf.add_paragraph()
    p_c.text = "Adaptive Context Engine"
    p_c.alignment = PP_ALIGN.CENTER
    p_c.font.name = FONT_BODY
    p_c.font.size = Pt(11)
    p_c.font.color.rgb = C_TEXT_MUTED
    p_c.space_before = Pt(2)

    p_d = h_tf.add_paragraph()
    p_d.text = (
        "Rather than an unguided chatbot, HumSaathi acts as an intelligent mediator. "
        "It ingests multi-dimensional contextual signals to dynamically calibrate dialogue pacing, "
        "cognitive challenge, and therapeutic reinforcement."
    )
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(10.5)
    p_d.font.color.rgb = C_TEXT_LIGHT
    p_d.space_before = Pt(14)

    # Left Context Inputs (Flowing into Hub)
    inputs_left = [
        ("LEARNER PROFILE", "Baseline sensory & diagnostic thresholds", C_EMERALD, Inches(2.2)),
        ("PERSONA CONTEXT", "Child · Teen · Adult life constraints", C_INDIGO, Inches(3.5)),
        ("SCENARIO & ROLE", "27 realistic social/workplace contexts", C_CYAN, Inches(4.8)),
    ]
    for tag, desc, col, top_y in inputs_left:
        ibox = s3.shapes.add_textbox(Inches(0.8), top_y, Inches(3.5), Inches(1.1))
        itf = ibox.text_frame
        itf.word_wrap = True
        p_tag = itf.paragraphs[0]
        p_tag.text = f"→  {tag}"
        p_tag.font.name = FONT_HEAD
        p_tag.font.size = Pt(12)
        p_tag.font.bold = True
        p_tag.font.color.rgb = col
        p_desc = itf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(2)

    # Right Context Outputs (Generated by Hub)
    inputs_right = [
        ("ADAPTIVE PACING", "Low-sensory calibration matching cognitive load", C_TEAL, Inches(2.2)),
        ("DYNAMIC ROLEPLAY", "Empathetic, realistic conversational simulation", C_CYAN, Inches(3.5)),
        ("OBJECTIVE SCORING", "Continuous evaluation: Clarity, Relevance, Confidence", C_AMBER, Inches(4.8)),
    ]
    for tag, desc, col, top_y in inputs_right:
        ibox = s3.shapes.add_textbox(Inches(9.0), top_y, Inches(3.5), Inches(1.1))
        itf = ibox.text_frame
        itf.word_wrap = True
        p_tag = itf.paragraphs[0]
        p_tag.text = f"{tag}  →"
        p_tag.font.name = FONT_HEAD
        p_tag.font.size = Pt(12)
        p_tag.font.bold = True
        p_tag.font.color.rgb = col
        p_desc = itf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(2)

    add_minimal_footer(s3, 3)

    s3.notes_slide.notes_text_frame.text = (
        "SLIDE 3 — THE BIG IDEA (Meet HumSaathi AI)\n"
        "TIME: 30 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"HumSaathi AI does not simply query a large language model and hope for the best.\n"
        "It is a centralized Context & Persona Engine. It ingests the learner's developmental stage, sensory preferences, "
        "and chosen scenario — whether it's navigating peer conflict or interviewing for a job.\n"
        "It then generates an adaptive, judgment-free learning experience with real-time conversational calibration.\"\n\n"
        "KEY JUDGE TAKEAWAY: HumSaathi is a full context-aware coaching engine, not a shallow OpenAI wrapper."
    )

    # ==========================================================================
    # SLIDE 04: THREE EXPERIENCES (Same Intelligence. Different Experience.)
    # ==========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_slide_base(s4, C_BG_DARK, is_morph=True)
    add_headline(s4, "Lifespan Adaptability", "Same Intelligence. Different Experience.")

    # 3 Floating Screen Mockups showing Child, Teen, Adult
    p_widths = Inches(3.7)
    p_heights = Inches(2.8)
    p_tops = Inches(2.2)

    # 1. Child
    add_browser_mockup(s4, shots["child"], Inches(0.8), p_tops, p_widths, p_heights, "humsaathi.ai/child")
    c_lbl = s4.shapes.add_textbox(Inches(0.8), p_tops + p_heights + Inches(0.4), p_widths, Inches(1.3))
    c_tf = c_lbl.text_frame
    c_tf.word_wrap = True
    p = c_tf.paragraphs[0]
    p.text = "CHILD EXPERIENCE  [Ages 4–11]"
    p.font.name = FONT_HEAD
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD
    p_d = c_tf.add_paragraph()
    p_d.text = "Visual learning, high predictability, phonics, numbers, emotions, and zero-anxiety routines."
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = C_TEXT_MUTED
    p_d.space_before = Pt(2)

    # 2. Teen
    add_browser_mockup(s4, shots["teen"], Inches(4.8), p_tops, p_widths, p_heights, "humsaathi.ai/teen")
    t_lbl = s4.shapes.add_textbox(Inches(4.8), p_tops + p_heights + Inches(0.4), p_widths, Inches(1.3))
    t_tf = t_lbl.text_frame
    t_tf.word_wrap = True
    p = t_tf.paragraphs[0]
    p.text = "TEEN EXPERIENCE  [Ages 12–17]"
    p.font.name = FONT_HEAD
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_INDIGO
    p_d = t_tf.add_paragraph()
    p_d.text = "Social dynamics, peer disputes, daily communication challenges, and authentic streak metrics."
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = C_TEXT_MUTED
    p_d.space_before = Pt(2)

    # 3. Adult
    add_browser_mockup(s4, shots["adult"], Inches(8.8), p_tops, p_widths, p_heights, "humsaathi.ai/adult")
    a_lbl = s4.shapes.add_textbox(Inches(8.8), p_tops + p_heights + Inches(0.4), p_widths, Inches(1.3))
    a_tf = a_lbl.text_frame
    a_tf.word_wrap = True
    p = a_tf.paragraphs[0]
    p.text = "ADULT EXPERIENCE  [Ages 18+]"
    p.font.name = FONT_HEAD
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    p_d = a_tf.add_paragraph()
    p_d.text = "Career readiness, workplace goal tracking, doctor & bank navigation, and contingency problem-solving."
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(10)
    p_d.font.color.rgb = C_TEXT_MUTED
    p_d.space_before = Pt(2)

    add_minimal_footer(s4, 4)

    s4.notes_slide.notes_text_frame.text = (
        "SLIDE 4 — THREE EXPERIENCES (Same Intelligence. Different Experience.)\n"
        "TIME: 30 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Here you see our actual, live platform adapting across life stages.\n"
        "For children: calm, sensory-safe exploration.\n"
        "For teens: peer dynamics, social challenges, and vocabulary.\n"
        "For adults: workplace meetings, job interviews, banking, and transit autonomy.\n"
        "One unified architecture; three deeply respectful, age-tailored user interfaces.\"\n\n"
        "KEY JUDGE TAKEAWAY: The platform solves the entire lifespan, preventing users from aging out of support."
    )

    # ==========================================================================
    # SLIDE 05: CHILD EXPERIENCE (Full-Screen Product Showcase)
    # ==========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_slide_base(s5, C_BG_DARK)
    add_headline(s5, "Early Foundations", "Child Portal: Calm, Structured Exploration")

    # Large Hero Browser Mockup on Left
    mock_w = Inches(8.0)
    mock_h = Inches(4.6)
    add_browser_mockup(s5, shots["child"], Inches(0.8), Inches(1.9), mock_w, mock_h, "humsaathi.ai/child/dashboard")

    # Right: 3 Subtle Callouts (Pure typography and whitespace)
    right_box = s5.shapes.add_textbox(Inches(9.2), Inches(2.2), Inches(3.3), Inches(4.3))
    r_tf = right_box.text_frame
    r_tf.word_wrap = True

    callouts = [
        ("VISUAL LEARNING", "High-contrast visual cues with bilingual English & Urdu audio pronunciation for letters, numbers, and animals."),
        ("LOW SENSORY LOAD", "Zero flashing animations, zero countdown clocks, and zero punishing fail buzzers. Pure cognitive safety."),
        ("PREDICTABLE INTERACTION", "Structured daily routines and emotion identification matching occupational therapy standards."),
    ]
    for i, (ctitle, cdesc) in enumerate(callouts):
        p_c = r_tf.paragraphs[0] if i == 0 else r_tf.add_paragraph()
        p_c.text = ctitle
        p_c.font.name = FONT_HEAD
        p_c.font.size = Pt(13)
        p_c.font.bold = True
        p_c.font.color.rgb = C_EMERALD
        if i > 0:
            p_c.space_before = Pt(16)

        p_cd = r_tf.add_paragraph()
        p_cd.text = cdesc
        p_cd.font.name = FONT_BODY
        p_cd.font.size = Pt(10.5)
        p_cd.font.color.rgb = C_TEXT_MUTED
        p_cd.space_before = Pt(4)

    add_minimal_footer(s5, 5)

    s5.notes_slide.notes_text_frame.text = (
        "SLIDE 5 — CHILD EXPERIENCE (Calm, Structured Exploration)\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Look at the actual Child Portal screen on the left.\n"
        "Notice the intentional whitespace and calm color system. Neurodiverse children easily experience sensory overload. "
        "We implemented 7 foundational categories: Letters, Numbers, Shapes, Colors, Animals, Emotions, and Daily Routines.\n"
        "Every interaction is predictable, bilingual, and positive.\"\n\n"
        "KEY JUDGE TAKEAWAY: Designed strictly with occupational therapy ergonomics in mind."
    )

    # ==========================================================================
    # SLIDE 06: TEEN + ADULT (Cinematic Split Screen)
    # ==========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_slide_base(s6, C_BG_DARK)
    add_headline(s6, "Maturity Progression", "Teen + Adult: Social Dynamics & Career Autonomy")

    split_w = Inches(5.65)
    split_h = Inches(3.6)

    # Left: TEEN
    add_browser_mockup(s6, shots["teen"], Inches(0.8), Inches(1.9), split_w, split_h, "humsaathi.ai/teen")
    t_desc = s6.shapes.add_textbox(Inches(0.8), Inches(5.8), split_w, Inches(1.1))
    td_tf = t_desc.text_frame
    td_tf.word_wrap = True
    p = td_tf.paragraphs[0]
    p.text = "TEEN: ACADEMIC & PEER DYNAMICS"
    p.font.name = FONT_HEAD
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_INDIGO
    p_b = td_tf.add_paragraph()
    p_b.text = "Features rotating Daily Communication Challenges (asking for help, expressing opinions) with real database-calculated contiguous streaks."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(10.5)
    p_b.font.color.rgb = C_TEXT_MUTED
    p_b.space_before = Pt(3)

    # Right: ADULT
    add_browser_mockup(s6, shots["adult"], Inches(6.88), Inches(1.9), split_w, split_h, "humsaathi.ai/adult")
    a_desc = s6.shapes.add_textbox(Inches(6.88), Inches(5.8), split_w, Inches(1.1))
    ad_tf = a_desc.text_frame
    ad_tf.word_wrap = True
    p = ad_tf.paragraphs[0]
    p.text = "ADULT: CAREER READINESS & INDEPENDENCE"
    p.font.name = FONT_HEAD
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    p_b = ad_tf.add_paragraph()
    p_b.text = "Workplace Goal tracking ('Complete 2 workplace scenarios'), manager check-ins, job interviews, banking, and transit delay problem-solving."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(10.5)
    p_b.font.color.rgb = C_TEXT_MUTED
    p_b.space_before = Pt(3)

    add_minimal_footer(s6, 6)

    s6.notes_slide.notes_text_frame.text = (
        "SLIDE 6 — TEEN + ADULT (Social Dynamics & Career Autonomy)\n"
        "TIME: 35 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Here is our side-by-side progression.\n"
        "On the left, the Teen Portal addresses social nuance: resolving group project disagreements, reading context clues, "
        "and building daily habits with genuine streak metrics.\n"
        "On the right, the Adult Portal transforms into a sleek, professional career hub. Adults prepare for job interviews, "
        "clarify ambiguous emails, and practice high-stress autonomy tasks like resolving bank debit card issues. Real tools for real independence.\"\n\n"
        "KEY JUDGE TAKEAWAY: Age-appropriate respect and dignity built into every pixel."
    )

    # ==========================================================================
    # SLIDE 07: HERO FEATURE — AI COMMUNICATION COACH (Live Conversation Engine)
    # ==========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_slide_base(s7, C_BG_HERO, is_morph=True)
    add_headline(s7, "Hero Capability", "AI That Doesn't Just Answer — It Practices With You")

    # Left: Conversation Interface Mockup
    conv_w = Inches(6.8)
    conv_h = Inches(4.6)
    add_browser_mockup(s7, shots["conversation"], Inches(0.8), Inches(1.9), conv_w, conv_h, "humsaathi.ai/conversation/session-live")

    # Right: The Diagnostic Intelligence Breakdown (Clean, minimal)
    intel_box = s7.shapes.add_textbox(Inches(8.0), Inches(1.9), Inches(4.5), Inches(4.8))
    i_tf = intel_box.text_frame
    i_tf.word_wrap = True

    p = i_tf.paragraphs[0]
    p.text = "MULTI-DIMENSIONAL DIAGNOSTIC"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    # Metrics
    m_data = [
        ("CLARITY SCORE", "85%", "Vocabulary precision and grammatical coherence"),
        ("RELEVANCE SCORE", "90%", "Directness in addressing situational prompt"),
        ("CONFIDENCE INDEX", "80%", "Conversational flow and reduced hesitation"),
    ]
    for m_t, m_v, m_d in m_data:
        p_m = i_tf.add_paragraph()
        p_m.text = f"{m_t}  —  {m_v}"
        p_m.font.name = FONT_HEAD
        p_m.font.size = Pt(13)
        p_m.font.bold = True
        p_m.font.color.rgb = C_TEXT_WHITE
        p_m.space_before = Pt(8)

        p_md = i_tf.add_paragraph()
        p_md.text = m_d
        p_md.font.name = FONT_BODY
        p_md.font.size = Pt(10)
        p_md.font.color.rgb = C_TEXT_MUTED

    # Highlight: Suggested Alternative Response
    p_alt_t = i_tf.add_paragraph()
    p_alt_t.text = "SUGGESTED ALTERNATIVE PHRASING"
    p_alt_t.font.name = FONT_HEAD
    p_alt_t.font.size = Pt(11)
    p_alt_t.font.bold = True
    p_alt_t.font.color.rgb = C_EMERALD
    p_alt_t.space_before = Pt(14)

    p_alt_q = i_tf.add_paragraph()
    p_alt_q.text = "“I understand your concern. The slides I took have the technical diagrams, but let's rebalance them right now so we both have an equal count.”"
    p_alt_q.font.name = FONT_BODY
    p_alt_q.font.size = Pt(10.5)
    p_alt_q.font.italic = True
    p_alt_q.font.color.rgb = C_TEXT_LIGHT
    p_alt_q.space_before = Pt(4)

    add_minimal_footer(s7, 7)

    s7.notes_slide.notes_text_frame.text = (
        "SLIDE 7 — AI COMMUNICATION COACH (Practices With You)\n"
        "TIME: 45 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"This is our core breakthrough.\n"
        "Generic chatbots write essays for you. HumSaathi AI roleplays *with* you.\n"
        "When a user says: 'I noticed the slides were distributed unevenly,' the AI acts as the teammate or manager, responding in character.\n"
        "At the end of the exchange, our engine scores Clarity, Relevance, and Confidence. "
        "And look at the bottom right: it models a 'Suggested Alternative Response' — showing the learner exactly how to rephrase "
        "their thought more constructively. That is real communication coaching.\"\n\n"
        "KEY JUDGE TAKEAWAY: Turn-by-turn simulation + structured evaluation + alternative phrasing model."
    )

    # ==========================================================================
    # SLIDE 08: TECHNICAL ARCHITECTURE (Cinematic Data-Flow Pipeline)
    # ==========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_slide_base(s8, C_BG_HERO)
    add_headline(s8, "Engineering Depth", "Production AI Architecture & Data Flow")

    # Horizontal Flow Nodes (Pure typography and connector aesthetics)
    nodes = [
        ("01. CLIENT LAYER", "React 18 + Vite SPA", "Web Speech API audio capture & I18n trilingual translation runtime", C_EMERALD),
        ("02. GATEWAY", "FastAPI ASGI", "Python 3.11 serverless routing with JWT bearer authentication", C_CYAN),
        ("03. CONTEXT ENGINE", "context_builder.py", "Referent anchor tracking & intent classification across dialogue turns", C_INDIGO),
        ("04. LLM & FALLBACK", "Gemini 2.0 / OpenAI", "Temperature-calibrated persona prompting with deterministic rule fallback", C_TEAL),
        ("05. EVALUATION", "evaluation_service.py", "Clarity, relevance, confidence scoring + alternative phrasing generation", C_AMBER),
        ("06. PERSISTENCE", "SQLAlchemy + DB", "Persona-isolated schema, streak computation & parent digest aggregation", C_EMERALD),
    ]

    node_w = Inches(1.85)
    node_h = Inches(4.3)
    for i, (layer, tech, detail, col) in enumerate(nodes):
        nx = Inches(0.8) + i * Inches(1.98)
        ny = Inches(2.1)

        box = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, nx, ny, node_w, node_h)
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE
        box.line.color.rgb = col
        box.line.width = Pt(1)

        tb = s8.shapes.add_textbox(nx + Inches(0.1), ny + Inches(0.15), node_w - Inches(0.2), node_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        p_l = tf.paragraphs[0]
        p_l.text = layer
        p_l.font.name = FONT_BODY
        p_l.font.size = Pt(8.5)
        p_l.font.bold = True
        p_l.font.color.rgb = col

        p_t = tf.add_paragraph()
        p_t.text = tech
        p_t.font.name = FONT_HEAD
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_WHITE
        p_t.space_before = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = detail
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(8)

    # Bottom Callout Banner
    bot_box = s8.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.733), Inches(0.4))
    p_b = bot_box.text_frame.paragraphs[0]
    p_b.text = "ENGINEERING HIGHLIGHT: Referent anchor tracking in context_builder.py resolves pronoun antecedents across 4+ turns, preventing topic drift."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(10)
    p_b.font.color.rgb = C_CYAN

    add_minimal_footer(s8, 8)

    s8.notes_slide.notes_text_frame.text = (
        "SLIDE 8 — PRODUCTION AI ARCHITECTURE (How The AI Works)\n"
        "TIME: 40 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"For our technical judges: how does HumSaathi actually work?\n"
        "From the React client, requests pass through our FastAPI ASGI gateway. "
        "Before calling any LLM, our Context and Persona Engine in context_builder.py resolves pronouns like 'what next?' or 'how so?', "
        "grounding the conversation to prevent topic drift.\n"
        "We prompt Gemini 2.0 or OpenAI under strict temperature calibration. If an API rate limit or outage occurs, "
        "our deterministic rules engine steps in automatically with zero downtime.\"\n\n"
        "KEY JUDGE TAKEAWAY: Proof of architectural sophistication, multi-turn statefulness, and production resilience."
    )

    # ==========================================================================
    # SLIDE 09: ADAPTIVE LEARNING LOOP (Dynamic Intelligence Cycle)
    # ==========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_slide_base(s9, C_BG_DARK, is_morph=True)
    add_headline(s9, "Adaptive Intelligence", "Continuous Calibration: Every Interaction Makes the Next Smarter")

    # 6 Steps arranged in a visual progression
    loop = [
        ("01. ASSESS", "Lightweight baseline diagnostic evaluating strengths and sensory preferences.", C_EMERALD),
        ("02. PRACTICE", "Engaging in interactive voice/text scenarios and foundational skill modules.", C_CYAN),
        ("03. OBSERVE", "Tracking response latency, sentence complexity, and vocabulary retention in real time.", C_INDIGO),
        ("04. EVALUATE", "Scoring Clarity, Relevance, and Confidence immediately post-session.", C_TEAL),
        ("05. ADAPT", "Recalibrating scenario difficulty thresholds and sensory stimulation levels.", C_AMBER),
        ("06. RECOMMEND", "Generating the 'Next Best Practice' card directly on the learner's dashboard.", C_EMERALD),
    ]

    for i, (title, desc, col) in enumerate(loop):
        col_x = Inches(0.8) + (i % 3) * Inches(3.99)
        row_y = Inches(2.2) + (i // 3) * Inches(2.3)

        box = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, row_y, Inches(3.7), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE
        box.line.color.rgb = col
        box.line.width = Pt(1)

        tb = s9.shapes.add_textbox(col_x + Inches(0.2), row_y + Inches(0.18), Inches(3.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = FONT_HEAD
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(6)

    add_minimal_footer(s9, 9)

    s9.notes_slide.notes_text_frame.text = (
        "SLIDE 9 — ADAPTIVE LEARNING LOOP\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"The core intelligence of HumSaathi is continuous calibration.\n"
        "Assess, Practice, Observe, Evaluate, Adapt, and Recommend.\n"
        "When an adult learner struggles with expressing a disagreement constructively, the engine detects hesitation and "
        "adapts: it lowers the difficulty tier and recommends an intermediate roleplay on active listening. "
        "The system grows with the learner.\"\n\n"
        "KEY JUDGE TAKEAWAY: Demonstrates the recommendation algorithm and personalized adaptation logic."
    )

    # ==========================================================================
    # SLIDE 10: VOICE + LANGUAGE (Minimal, Impactful, Trilingual)
    # ==========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_slide_base(s10, C_BG_HERO)
    add_headline(s10, "Multimodal Accessibility", "Speak. Listen. Practice in Your Language.")

    # Massive Typography Statement
    big_t = s10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(2.5))
    b_tf = big_t.text_frame
    b_tf.word_wrap = True
    p1 = b_tf.paragraphs[0]
    p1.text = "SPEAK.\nLISTEN.\nPRACTICE."
    p1.font.name = FONT_HEAD
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = C_TEXT_WHITE

    sub_t = s10.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(5.8), Inches(1.8))
    st_tf = sub_t.text_frame
    st_tf.word_wrap = True
    p = st_tf.paragraphs[0]
    p.text = "Web Speech API integration powers authentic real-world phone calls and conversational speech recognition without third-party vendor latency."
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = C_TEXT_MUTED

    # Right: 3 Language Identity Blocks
    langs = [
        ("ENGLISH", "en", "Standard international terminology for academic and professional workplace mastery.", C_CYAN),
        ("URDU  (اردو)", "ur", "Full native Right-to-Left (RTL) typography with culturally grounded phrasing.", C_EMERALD),
        ("ROMAN URDU", "ur_rm", "Phonetic conversational vernacular widely spoken across South Asian diaspora communities.", C_INDIGO),
    ]
    for i, (name, code, desc, col) in enumerate(langs):
        ly = Inches(2.0) + i * Inches(1.6)
        box = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), ly, Inches(5.3), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE
        box.line.color.rgb = col
        box.line.width = Pt(1)

        tb = s10.shapes.add_textbox(Inches(7.4), ly + Inches(0.15), Inches(4.9), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = f"{name}  [{code}]"
        p_t.font.name = FONT_HEAD
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_LIGHT
        p_d.space_before = Pt(3)

    add_minimal_footer(s10, 10)

    s10.notes_slide.notes_text_frame.text = (
        "SLIDE 10 — VOICE + LANGUAGE (Speak. Listen. Practice.)\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Real communication is spoken. Using the browser's Web Speech API, learners practice phone calls and face-to-face dialogue.\n"
        "And we are natively trilingual: English, native Urdu with full RTL typography, and Roman Urdu. "
        "By supporting Roman Urdu, we break down the literacy barrier for hundreds of millions in South Asia and global diasporas.\"\n\n"
        "KEY JUDGE TAKEAWAY: True multimodal speech recognition + cultural localization."
    )

    # ==========================================================================
    # SLIDE 11: PARENT EXPERIENCE (Intelligent Caregiver Analytics)
    # ==========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_slide_base(s11, C_BG_DARK)
    add_headline(s11, "Family Ecosystem", "Parent Companion: Understand Progress — Not Control It")

    # Left: Actual Parent Portal Screenshot
    add_browser_mockup(s11, shots["parent"], Inches(0.8), Inches(1.9), Inches(7.2), Inches(4.7), "humsaathi.ai/parent/communication")

    # Right: The 4 Communication Dimensions
    r_box = s11.shapes.add_textbox(Inches(8.4), Inches(2.0), Inches(4.1), Inches(4.7))
    r_tf = r_box.text_frame
    r_tf.word_wrap = True

    p = r_tf.paragraphs[0]
    p.text = "4 CORE COMMUNICATION DIMENSIONS"
    p.font.name = FONT_BODY
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD

    dims = [
        ("GREETING & INITIATION", "Measures the learner's ability to initiate conversations comfortably."),
        ("ANSWERING QUESTIONS", "Evaluates directness, relevance, and response accuracy."),
        ("ASKING QUESTIONS", "Tracks proactive inquiry, requesting help, and conversational curiosity."),
        ("CONFIDENCE INDEX", "Monitors conversational flow, pacing, and reduced hesitation."),
    ]
    for dtitle, ddesc in dims:
        p_t = r_tf.add_paragraph()
        p_t.text = dtitle
        p_t.font.name = FONT_HEAD
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_WHITE
        p_t.space_before = Pt(8)

        p_d = r_tf.add_paragraph()
        p_d.text = ddesc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(2)

    # Privacy notice
    p_priv = r_tf.add_paragraph()
    p_priv.text = "🔒 PRIVACY-FIRST: Parents receive AI weekly digests and home practice guidance without reading private transcripts."
    p_priv.font.name = FONT_BODY
    p_priv.font.size = Pt(9.5)
    p_priv.font.color.rgb = C_CYAN
    p_priv.space_before = Pt(12)

    add_minimal_footer(s11, 11)

    s11.notes_slide.notes_text_frame.text = (
        "SLIDE 11 — PARENT COMPANION (Understand Progress — Not Control It)\n"
        "TIME: 30 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Look at the actual Parent Portal screen.\n"
        "Instead of leaving parents wondering how their child is doing, we track four core communication dimensions: "
        "Greeting, Answering, Asking, and Confidence.\n"
        "And we respect learner dignity: parents do not read private chat transcripts. Our AI synthesizes home guidance, "
        "like 'Encourage Zayd to order dinner tonight to reinforce his progress in asking questions.' Protected behind a 4-digit PIN.\"\n\n"
        "KEY JUDGE TAKEAWAY: Actionable caregiver insights without invasive surveillance."
    )

    # ==========================================================================
    # SLIDE 12: ENGINEERING / TECHNOLOGY (Clean Layered Visual)
    # ==========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_slide_base(s12, C_BG_DARK)
    add_headline(s12, "Production Engineering", "Full-Stack AI Ecosystem")

    layers = [
        ("FRONTEND", "React 18  ·  Vite  ·  React Router v7  ·  Native CSS Design Tokens", "Sub-second initial paint, zero bloated component libraries, full accessibility contrast.", C_CYAN),
        ("BACKEND API", "Python 3.11  ·  FastAPI  ·  Pydantic v2  ·  Uvicorn ASGI", "Asynchronous serverless endpoints with microsecond request validation and rate limiting.", C_EMERALD),
        ("AI ENGINE", "Google Gemini 2.0  ·  OpenAI GPT-4o  ·  Deterministic Rule Engine", "Referent anchor resolution, intent classification, and instant offline fallback resiliency.", C_INDIGO),
        ("PERSISTENCE", "SQLAlchemy ORM  ·  SQLite (Local)  ·  PostgreSQL (Cloud)", "Relational data schema with strict persona isolation filters and real-time streak computation.", C_TEAL),
        ("INFRASTRUCTURE", "Vercel Serverless Functions  ·  Global Edge CDN  ·  GitHub CI/CD", "Automated deployment pipeline with 100% production uptime at hum-saathi-ai.vercel.app.", C_CYAN),
    ]

    for i, (layer_name, tech_name, desc, col) in enumerate(layers):
        ly = Inches(2.0) + i * Inches(0.95)
        box = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), ly, Inches(11.733), Inches(0.82))
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE
        box.line.color.rgb = C_BORDER
        box.line.width = Pt(1)

        # Left tag pill
        tag_b = s12.shapes.add_textbox(Inches(1.0), ly + Inches(0.18), Inches(1.8), Inches(0.45))
        p_tag = tag_b.text_frame.paragraphs[0]
        p_tag.text = layer_name
        p_tag.font.name = FONT_HEAD
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = col

        # Middle tech name
        tech_b = s12.shapes.add_textbox(Inches(2.8), ly + Inches(0.18), Inches(4.5), Inches(0.45))
        p_tech = tech_b.text_frame.paragraphs[0]
        p_tech.text = tech_name
        p_tech.font.name = FONT_BODY
        p_tech.font.size = Pt(11)
        p_tech.font.bold = True
        p_tech.font.color.rgb = C_TEXT_WHITE

        # Right desc
        desc_b = s12.shapes.add_textbox(Inches(7.4), ly + Inches(0.18), Inches(4.9), Inches(0.45))
        p_desc = desc_b.text_frame.paragraphs[0]
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = C_TEXT_MUTED

    add_minimal_footer(s12, 12)

    s12.notes_slide.notes_text_frame.text = (
        "SLIDE 12 — ENGINEERING / TECHNOLOGY\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Here is our actual production stack.\n"
        "React 18 and Vite for instant load times.\n"
        "Python 3.11 with FastAPI for serverless ASGI performance.\n"
        "Google Gemini 2.0 paired with our deterministic offline fallback engine.\n"
        "Everything is version-controlled and deployed live on Vercel's edge network right now.\"\n\n"
        "KEY JUDGE TAKEAWAY: Modern, production-ready, clean engineering."
    )

    # ==========================================================================
    # SLIDE 13: SECURITY (Secure Architecture Data Flow)
    # ==========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_slide_base(s13, C_BG_DARK)
    add_headline(s13, "Data Protection", "Security & Privacy Built Into the Architecture")

    # Flow sequence: IDENTITY -> ROLE -> PERSONA -> DATA -> ACCESS
    flow_steps = [
        ("01. IDENTITY", "JWT & Google OAuth 2.0", "Cryptographically verified authentication with passwordless login support."),
        ("02. ROLE CONTROL", "Enterprise RBAC", "Strict separation of Learner, Parent, and System Administrator privileges."),
        ("03. PERSONA ISOLATION", "Database Partitioning", "Adult and Child progress metrics are query-isolated, preventing data leakage."),
        ("04. CAREGIVER GATE", "Encrypted 4-Digit PIN", "Parent portal access locked locally, preventing accidental profile edits."),
    ]

    for i, (stitle, ssub, sdesc) in enumerate(flow_steps):
        sx = Inches(0.8) + i * Inches(2.98)
        box = s13.shapes.add_shape(MSO_SHAPE.RECTANGLE, sx, Inches(2.2), Inches(2.78), Inches(4.3))
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE
        box.line.color.rgb = C_BORDER
        box.line.width = Pt(1)

        tb = s13.shapes.add_textbox(sx + Inches(0.18), Inches(2.4), Inches(2.42), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = stitle
        p_t.font.name = FONT_BODY
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = C_CYAN

        p_s = tf.add_paragraph()
        p_s.text = ssub
        p_s.font.name = FONT_HEAD
        p_s.font.size = Pt(13)
        p_s.font.bold = True
        p_s.font.color.rgb = C_TEXT_WHITE
        p_s.space_before = Pt(4)

        p_d = tf.add_paragraph()
        p_d.text = sdesc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(10.5)
        p_d.font.color.rgb = C_TEXT_MUTED
        p_d.space_before = Pt(8)

    add_minimal_footer(s13, 13)

    s13.notes_slide.notes_text_frame.text = (
        "SLIDE 13 — SECURITY (Privacy Built Into the Architecture)\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"When dealing with neurodiverse and younger learners, security is non-negotiable.\n"
        "We enforce database-level persona data isolation — child records and adult workplace scenarios cannot cross-contaminate.\n"
        "Parent access is protected by an encrypted PIN, and our admin panel provides full audit logging and AI monitoring.\"\n\n"
        "KEY JUDGE TAKEAWAY: Enterprise-grade compliance and ethical AI safeguards."
    )

    # ==========================================================================
    # SLIDE 14: WHY HUMSAATHI? (Dramatic Visual Transformation)
    # ==========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_slide_base(s14, C_BG_HERO, is_morph=True)
    add_headline(s14, "The Value Proposition", "Why HumSaathi AI? More Than an AI Chatbot")

    # Split: Traditional vs HumSaathi (Typography-driven transformation)
    # Left: TRADITIONAL
    t_box = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.7))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = C_SURFACE
    t_box.line.color.rgb = C_ROSE
    t_box.line.width = Pt(1)

    t_tb = s14.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(5.0), Inches(4.1))
    t_tf = t_tb.text_frame
    t_tf.word_wrap = True

    p = t_tf.paragraphs[0]
    p.text = "TRADITIONAL PARADIGMS"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_ROSE

    t_points = [
        ("Generic LLM Prompt Box", "Overwhelming for neurodiverse learners; causes executive dysfunction and cognitive shutdown."),
        ("Static Flashcards & Quizzes", "Teaches passive vocabulary memorization with zero conversational agility."),
        ("Binary Right / Wrong Marks", "Provides pass/fail scores with no actionable guidance on tone, social nuance, or clarity."),
        ("One Audience Assumption", "Patronizes adults with child cartoons or overwhelms children with dense menus."),
    ]
    for pt, pd in t_points:
        p1 = t_tf.add_paragraph()
        p1.text = f"✕  {pt}"
        p1.font.name = FONT_HEAD
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_WHITE
        p1.space_before = Pt(6)
        p2 = t_tf.add_paragraph()
        p2.text = pd
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_MUTED

    # Right: HUMSAATHI AI
    h_box = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.7))
    h_box.fill.solid()
    h_box.fill.fore_color.rgb = C_SURFACE_HI
    h_box.line.color.rgb = C_EMERALD
    h_box.line.width = Pt(2)

    h_tb = s14.shapes.add_textbox(Inches(7.2), Inches(2.3), Inches(5.0), Inches(4.1))
    h_tf = h_tb.text_frame
    h_tf.word_wrap = True

    p = h_tf.paragraphs[0]
    p.text = "HUMSAATHI AI INNOVATION"
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD

    h_points = [
        ("Scaffolded Scenario Roleplay", "27 goal-oriented real-world situations with empathetic turn-by-turn guidance."),
        ("Multidimensional Evaluation", "Objective scoring of Clarity, Relevance, and Confidence post-session."),
        ("Suggested Alternative Phrasing", "Teaches learners how to rephrase thoughts more constructively next time."),
        ("Lifespan Personalization", "Deeply tailored Child, Teen, and Adult portals with parent observation insights."),
    ]
    for pt, pd in h_points:
        p1 = h_tf.add_paragraph()
        p1.text = f"✓  {pt}"
        p1.font.name = FONT_HEAD
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_EMERALD
        p1.space_before = Pt(6)
        p2 = h_tf.add_paragraph()
        p2.text = pd
        p2.font.name = FONT_BODY
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_LIGHT

    add_minimal_footer(s14, 14)

    s14.notes_slide.notes_text_frame.text = (
        "SLIDE 14 — WHY HUMSAATHI? (The Value Proposition)\n"
        "TIME: 30 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"Judges often ask: 'Why not just use ChatGPT?'\n"
        "Here is the answer: an open text box is paralyzing for a neurodiverse learner.\n"
        "HumSaathi provides scaffolded structure. It evaluates Clarity and Confidence. It offers direct alternative phrasing. "
        "And it adapts to the learner's developmental stage.\n"
        "It is a specialized communication health intervention, not a text toy.\"\n\n"
        "KEY JUDGE TAKEAWAY: Clear competitive moat established."
    )

    # ==========================================================================
    # SLIDE 15: CINEMATIC ENDING (Minimal, Memorable, Product Launch Finale)
    # ==========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_slide_base(s15, C_BG_HERO, is_morph=True)

    # Three core statements in massive typography
    core_box = s15.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(10.933), Inches(2.6))
    c_tf = core_box.text_frame
    c_tf.word_wrap = True

    p1 = c_tf.paragraphs[0]
    p1.text = "LEARNING IS PERSONAL."
    p1.font.name = FONT_HEAD
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = C_TEXT_DIM

    p2 = c_tf.add_paragraph()
    p2.text = "COMMUNICATION IS PRACTICE."
    p2.font.name = FONT_HEAD
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = C_CYAN
    p2.space_before = Pt(4)

    p3 = c_tf.add_paragraph()
    p3.text = "GROWTH IS CONTINUOUS."
    p3.font.name = FONT_HEAD
    p3.font.size = Pt(32)
    p3.font.bold = True
    p3.font.color.rgb = C_TEXT_DIM
    p3.space_before = Pt(4)

    # Horizontal hairline
    div = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.2), Inches(10.933), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = C_BORDER
    div.line.fill.background()

    # Final Hero Lockup
    bot_b = s15.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.933), Inches(2.2))
    b_tf = bot_b.text_frame
    b_tf.word_wrap = True

    p_h = b_tf.paragraphs[0]
    p_h.text = "HUMSAATHI AI"
    p_h.font.name = FONT_HEAD
    p_h.font.size = Pt(42)
    p_h.font.bold = True
    p_h.font.color.rgb = C_TEXT_WHITE

    p_tag = b_tf.add_paragraph()
    p_tag.text = "Learn. Communicate. Grow."
    p_tag.font.name = FONT_BODY
    p_tag.font.size = Pt(18)
    p_tag.font.bold = True
    p_tag.font.color.rgb = C_CYAN
    p_tag.space_before = Pt(2)

    p_quote = b_tf.add_paragraph()
    p_quote.text = "“HumSaathi AI doesn't just teach users what to learn — it helps them learn how to communicate, practice, and grow.”"
    p_quote.font.name = FONT_BODY
    p_quote.font.size = Pt(14)
    p_quote.font.italic = True
    p_quote.font.color.rgb = C_TEXT_LIGHT
    p_quote.space_before = Pt(8)

    add_minimal_footer(s15, 15)

    s15.notes_slide.notes_text_frame.text = (
        "SLIDE 15 — CINEMATIC ENDING (Learn. Communicate. Grow.)\n"
        "TIME: 25 Seconds\n\n"
        "TALKING POINTS:\n"
        "\"To conclude:\n"
        "Learning is personal.\n"
        "Communication is practice.\n"
        "Growth is continuous.\n\n"
        "HumSaathi AI doesn't just teach users what to learn — it helps them learn how to communicate, practice, and grow.\n"
        "The platform is deployed live at hum-saathi-ai.vercel.app. We invite you to experience it, and we welcome your questions!\"\n\n"
        "KEY JUDGE TAKEAWAY: A cinematic, confident, and unforgettable finish."
    )

    # Save to HumSaathi_AI_Hackathon_Pitch.pptx
    out_file = "HumSaathi_AI_Hackathon_Pitch.pptx"
    prs.save(out_file)
    print(f"SUCCESS: Redesigned and generated {out_file} with 15 cinematic slides and real screenshots!")

if __name__ == "__main__":
    build_cinematic_deck()
