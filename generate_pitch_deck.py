import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

# ==============================================================================
# COLOR PALETTE DEFINITION (Sophisticated Navy, Slate, Cyan, Emerald & Indigo)
# ==============================================================================
BG_DARK = RGBColor(11, 19, 43)        # #0B132B Deep rich navy
BG_DARKER = RGBColor(7, 12, 28)       # #070C1C Cinematic deep navy
CARD_DARK = RGBColor(21, 32, 53)      # #152035 Sleek container
CARD_ALT = RGBColor(28, 42, 68)       # #1C2A44 Hover/accent container
BORDER_COLOR = RGBColor(46, 64, 94)   # #2E405E Subtle border line

ACCENT_CYAN = RGBColor(6, 182, 212)    # #06B6D4 Vibrant cyan
ACCENT_GREEN = RGBColor(16, 185, 129)  # #10B981 Emerald
ACCENT_INDIGO = RGBColor(99, 102, 241) # #6366F1 Modern indigo
ACCENT_PURPLE = RGBColor(168, 85, 247) # #A855F7 Violet
ACCENT_ORANGE = RGBColor(245, 158, 11) # #F59E0B Warm amber

TEXT_WHITE = RGBColor(248, 250, 252)  # #F8FAFC Pure crisp white
TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Secondary silver
TEXT_DIM = RGBColor(100, 116, 139)    # #64748B Tertiary slate

FONT_MAIN = "Segoe UI"
FONT_HEADING = "Segoe UI Semibold"

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    logo_path = os.path.abspath("frontend/public/humsaathi-logo-v1.png")
    has_logo = os.path.exists(logo_path)

    # --------------------------------------------------------------------------
    # HELPER FUNCTIONS
    # --------------------------------------------------------------------------
    def add_bg(slide, dark_color=BG_DARK):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = dark_color
        bg.line.fill.background()
        
        # Add transition XML for smooth slide presentation
        trans_xml = parse_xml(
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med">'
            '<p:fade/></p:transition>'
        )
        slide._element.append(trans_xml)
        return bg

    def add_header(slide, eyebrow, title, subtitle=None):
        # Eyebrow
        eb_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.4))
        eb_tf = eb_box.text_frame
        eb_tf.word_wrap = True
        eb_tf.margin_left = eb_tf.margin_top = eb_tf.margin_right = eb_tf.margin_bottom = 0
        p = eb_tf.paragraphs[0]
        p.text = eyebrow.upper()
        p.font.name = FONT_MAIN
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.8))
        t_tf = t_box.text_frame
        t_tf.word_wrap = True
        t_tf.margin_left = t_tf.margin_top = t_tf.margin_right = t_tf.margin_bottom = 0
        p2 = t_tf.paragraphs[0]
        p2.text = title
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.4))
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sub_tf.margin_left = sub_tf.margin_top = sub_tf.margin_right = sub_tf.margin_bottom = 0
            p3 = sub_tf.paragraphs[0]
            p3.text = subtitle
            p3.font.name = FONT_MAIN
            p3.font.size = Pt(14)
            p3.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=CARD_DARK, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
        return shape

    def add_footer(slide, slide_num):
        # Footer watermark
        f_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        f_tf = f_box.text_frame
        f_tf.word_wrap = True
        f_tf.margin_left = f_tf.margin_top = f_tf.margin_right = f_tf.margin_bottom = 0
        p = f_tf.paragraphs[0]
        p.text = f"HumSaathi AI  |  Adaptive Communication & Learning Platform  |  Slide {slide_num:02d}"
        p.font.name = FONT_MAIN
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_DIM

    # ==========================================================================
    # SLIDE 01: HERO
    # ==========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    add_bg(s1, BG_DARKER)

    # Glowing background card
    add_card(s1, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1), bg_color=CARD_DARK, border_color=BORDER_COLOR)

    # Logo image if available
    if has_logo:
        try:
            s1.shapes.add_picture(logo_path, Inches(5.9), Inches(1.6), width=Inches(1.5))
        except Exception:
            pass

    # Title Box
    h_box = s1.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(9.333), Inches(1.2))
    h_tf = h_box.text_frame
    h_tf.word_wrap = True
    p = h_tf.paragraphs[0]
    p.text = "HumSaathi AI"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADING
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # Tagline
    tag_box = s1.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.333), Inches(0.6))
    tag_tf = tag_box.text_frame
    p_tag = tag_tf.paragraphs[0]
    p_tag.text = "Learn. Communicate. Grow."
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.name = FONT_MAIN
    p_tag.font.size = Pt(22)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_CYAN

    # Subtitle
    sub_box = s1.shapes.add_textbox(Inches(2.0), Inches(4.75), Inches(9.333), Inches(0.5))
    sub_tf = sub_box.text_frame
    p_sub = sub_tf.paragraphs[0]
    p_sub.text = "Adaptive AI Learning & Communication Coach for Neurodiverse Learners"
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.name = FONT_MAIN
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = TEXT_MUTED

    # Persona Badges
    badge_data = [
        ("🧒 CHILD", "Ages 4–11", ACCENT_GREEN, Inches(3.4)),
        ("🧑 TEEN", "Ages 12–17", ACCENT_INDIGO, Inches(5.8)),
        ("👨 ADULT", "Ages 18+", ACCENT_CYAN, Inches(8.2)),
    ]
    for name, age, col, left_pos in badge_data:
        b_card = add_card(s1, left_pos, Inches(5.35), Inches(1.8), Inches(0.65), bg_color=CARD_ALT, border_color=col)
        b_box = s1.shapes.add_textbox(left_pos, Inches(5.38), Inches(1.8), Inches(0.6))
        b_tf = b_box.text_frame
        p_b1 = b_tf.paragraphs[0]
        p_b1.text = name
        p_b1.alignment = PP_ALIGN.CENTER
        p_b1.font.name = FONT_MAIN
        p_b1.font.size = Pt(11)
        p_b1.font.bold = True
        p_b1.font.color.rgb = TEXT_WHITE
        p_b2 = b_tf.add_paragraph()
        p_b2.text = age
        p_b2.alignment = PP_ALIGN.CENTER
        p_b2.font.name = FONT_MAIN
        p_b2.font.size = Pt(9.5)
        p_b2.font.color.rgb = col

    add_footer(s1, 1)

    s1.notes_slide.notes_text_frame.text = (
        "SLIDE 1 — HERO / TITLE (Learn. Communicate. Grow.)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Good morning, esteemed judges. Over 1 billion people worldwide navigate neurodiversity and speech-language "
        "challenges. Yet today's educational tools treat learning as a rigid, one-way street. "
        "We are proud to introduce HumSaathi AI — an adaptive communication and learning companion engineered "
        "to empower neurodiverse learners across their entire lifespan: Child, Teen, and Adult. "
        "HumSaathi doesn't just teach facts — it coaches real-world communication confidence.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Emphasize lifespan coverage (Child, Teen, Adult).\n"
        "- Frame communication as the core barrier being solved.\n\n"
        "DEMO TRANSITION:\n"
        "\"Before showing you how HumSaathi AI works, let's examine why traditional educational technology fails these learners.\""
    )

    # ==========================================================================
    # SLIDE 02: THE PROBLEM
    # ==========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_bg(s2)
    add_header(s2, "The Challenge", "Learning Shouldn't Be One-Size-Fits-All",
               "Traditional educational tools fail neurodiverse individuals by treating learning as a rigid, one-size-fits-all conveyor belt.")

    prob_cards = [
        ("STATIC CONTENT", "One-Way Curricula",
         "Standard apps present the identical linear lessons regardless of sensory sensitivities, processing pace, or cognitive fatigue.",
         ACCENT_ORANGE, "📦"),
        ("PRACTICE GAP", "Zero Safe Roleplay",
         "Learners can memorize vocabulary but have no safe, judgment-free space to practice real-world social and conversational exchanges.",
         ACCENT_PURPLE, "🗣️"),
        ("FEEDBACK VOID", "Binary Right/Wrong",
         "Traditional tools return simple pass/fail checks, offering no actionable coaching on tone, social appropriateness, or confidence.",
         ACCENT_CYAN, "📊"),
        ("ANXIETY BARRIER", "High Social Fatigue",
         "Real-world conversations trigger sensory overload and fear of judgment, preventing independent communication growth.",
         ACCENT_INDIGO, "⚡"),
    ]

    card_w = Inches(2.78)
    card_h = Inches(4.6)
    for i, (title, sub, body, col, icon) in enumerate(prob_cards):
        cx = Inches(0.8) + i * Inches(2.98)
        cy = Inches(2.1)
        add_card(s2, cx, cy, card_w, card_h, bg_color=CARD_DARK, border_color=BORDER_COLOR)

        # Top accent pill
        add_card(s2, cx + Inches(0.2), cy + Inches(0.25), Inches(0.6), Inches(0.6), bg_color=CARD_ALT, border_color=col)
        ic_box = s2.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.28), Inches(0.6), Inches(0.5))
        ic_tf = ic_box.text_frame
        p_ic = ic_tf.paragraphs[0]
        p_ic.text = icon
        p_ic.alignment = PP_ALIGN.CENTER
        p_ic.font.size = Pt(16)

        # Category
        cat_box = s2.shapes.add_textbox(cx + Inches(0.2), cy + Inches(1.1), Inches(2.38), Inches(0.3))
        cat_tf = cat_box.text_frame
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = title
        p_cat.font.name = FONT_MAIN
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = col

        # Subtitle
        sub_b = s2.shapes.add_textbox(cx + Inches(0.2), cy + Inches(1.4), Inches(2.38), Inches(0.6))
        sub_b_tf = sub_b.text_frame
        sub_b_tf.word_wrap = True
        p_s = sub_b_tf.paragraphs[0]
        p_s.text = sub
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(16)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE

        # Description
        desc_box = s2.shapes.add_textbox(cx + Inches(0.2), cy + Inches(2.1), Inches(2.38), Inches(2.2))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        p_d = desc_tf.paragraphs[0]
        p_d.text = body
        p_d.font.name = FONT_MAIN
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = TEXT_MUTED

    add_footer(s2, 2)

    s2.notes_slide.notes_text_frame.text = (
        "SLIDE 2 — THE PROBLEM (Learning Shouldn't Be One-Size-Fits-All)\n\n"
        "SPEAKING TIME: 35 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Why is communication learning broken for neurodiverse individuals? We identified four critical fractures:\n"
        "First, Static Content — apps force every learner through the same rigid lessons.\n"
        "Second, The Practice Gap — knowing a word isn't the same as knowing how to use it in a conversation.\n"
        "Third, The Feedback Void — red checkmarks don't teach conversational flow or social nuance.\n"
        "And fourth, High Anxiety — facing unpredictable real-world interactions without safe simulation causes withdrawal.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Show that this isn't just about general education; it's about social and emotional independence.\n\n"
        "TRANSITION:\n"
        "\"This is exactly why we built HumSaathi AI.\""
    )

    # ==========================================================================
    # SLIDE 03: THE BIG IDEA
    # ==========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_bg(s3)
    add_header(s3, "The Solution", "Meet HumSaathi AI",
               "An adaptive intelligence engine connecting personalized learning, real-world communication practice, and measurable progress.")

    # Flow Steps
    steps = [
        ("1. LEARNER", "Personal Profile", "Calibrated sensory & developmental baseline", ACCENT_GREEN, Inches(0.8)),
        ("2. AI COMPANION", "Context Engine", "Empathetic, persona-aware conversational grounding", ACCENT_CYAN, Inches(3.25)),
        ("3. ADAPTIVE CONTENT", "Micro-Modules", "Dynamic pacing matching cognitive capacity", ACCENT_INDIGO, Inches(5.7)),
        ("4. PRACTICE SCENARIOS", "Realistic Roleplay", "Workplace, peer & everyday voice/text dialogues", ACCENT_PURPLE, Inches(8.15)),
        ("5. GROWTH", "Real Progress", "Contiguous streak & parent observation reports", ACCENT_ORANGE, Inches(10.6)),
    ]

    for title, sub, body, col, lx in steps:
        add_card(s3, lx, Inches(2.2), Inches(2.0), Inches(3.8), bg_color=CARD_DARK, border_color=col)
        
        # Pill
        add_card(s3, lx + Inches(0.15), Inches(2.4), Inches(1.7), Inches(0.4), bg_color=CARD_ALT, border_color=col)
        p_box = s3.shapes.add_textbox(lx + Inches(0.15), Inches(2.45), Inches(1.7), Inches(0.35))
        p_tf = p_box.text_frame
        p_par = p_tf.paragraphs[0]
        p_par.text = title
        p_par.alignment = PP_ALIGN.CENTER
        p_par.font.size = Pt(10)
        p_par.font.bold = True
        p_par.font.color.rgb = col

        # Sub
        s_box = s3.shapes.add_textbox(lx + Inches(0.15), Inches(3.0), Inches(1.7), Inches(0.8))
        s_tf = s_box.text_frame
        s_tf.word_wrap = True
        p_s = s_tf.paragraphs[0]
        p_s.text = sub
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(15)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE

        # Body
        b_box = s3.shapes.add_textbox(lx + Inches(0.15), Inches(3.8), Inches(1.7), Inches(1.9))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True
        p_b = b_tf.paragraphs[0]
        p_b.text = body
        p_b.font.name = FONT_MAIN
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED

    # Banner Card underneath
    banner = add_card(s3, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.7), bg_color=CARD_ALT, border_color=ACCENT_CYAN)
    ban_box = s3.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.333), Inches(0.5))
    ban_tf = ban_box.text_frame
    p_ban = ban_tf.paragraphs[0]
    p_ban.text = "CORE MISSION: Transforming passive screen consumption into active, self-paced, judgment-free communication mastery."
    p_ban.alignment = PP_ALIGN.CENTER
    p_ban.font.name = FONT_HEADING
    p_ban.font.size = Pt(13)
    p_ban.font.bold = True
    p_ban.font.color.rgb = ACCENT_CYAN

    add_footer(s3, 3)

    s3.notes_slide.notes_text_frame.text = (
        "SLIDE 3 — THE BIG IDEA (Meet HumSaathi AI)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Meet HumSaathi AI — which means 'Our Companion'.\n"
        "Instead of asking the learner to conform to the software, HumSaathi conforms to the learner.\n"
        "The learner is paired with an AI Companion that personalizes content, orchestrates realistic roleplay, "
        "and observes communication patterns in real-time. The result is genuine, measurable growth.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Emphasize that HumSaathi is an interactive coach, not a static quiz engine."
    )

    # ==========================================================================
    # SLIDE 04: ONE PLATFORM. THREE EXPERIENCES.
    # ==========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_bg(s4)
    add_header(s4, "Lifespan Adaptability", "One Platform. Three Experiences.",
               "Neurodiversity doesn't stop at childhood. HumSaathi AI provides three deeply differentiated, age-tailored environments.")

    portals = [
        ("CHILD PORTAL", "Learn & Explore", "Ages 4 to 11", ACCENT_GREEN, [
            "Low sensory load & high predictable visual structure",
            "Foundational exploration: Letters, Numbers, Shapes & Colors",
            "Emotions & daily routine identification with gentle audio",
            "Encouraging, gamified positive reinforcement"
        ]),
        ("TEEN PORTAL", "Practice & Communicate", "Ages 12 to 17", ACCENT_INDIGO, [
            "Peer communication, social dynamics & empathy building",
            "Daily Communication Challenges (Saying 'I need help', opinions)",
            "Context clues, reading passages & social dispute resolution",
            "Real streak tracking with zero fake statistics"
        ]),
        ("ADULT PORTAL", "Apply & Grow", "Ages 18 and Above", ACCENT_CYAN, [
            "Workplace readiness: Job interviews & manager discussions",
            "Workplace Goal Engine: 'Complete 2 workplace scenarios'",
            "Everyday life autonomy: Doctor visits, banking, dining",
            "Practical problem solving: Transit delays & confusing emails"
        ]),
    ]

    card_pw = Inches(3.75)
    card_ph = Inches(4.7)
    for i, (title, sub, age, col, bullets) in enumerate(portals):
        px = Inches(0.8) + i * Inches(3.99)
        py = Inches(2.1)
        add_card(s4, px, py, card_pw, card_ph, bg_color=CARD_DARK, border_color=col)

        # Header Badge
        add_card(s4, px + Inches(0.25), py + Inches(0.3), Inches(3.25), Inches(0.55), bg_color=CARD_ALT, border_color=col)
        h_b = s4.shapes.add_textbox(px + Inches(0.25), py + Inches(0.35), Inches(3.25), Inches(0.45))
        h_tf = h_b.text_frame
        p_h = h_tf.paragraphs[0]
        p_h.text = f"{title} · {age}"
        p_h.alignment = PP_ALIGN.CENTER
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = col

        # Subtitle
        s_b = s4.shapes.add_textbox(px + Inches(0.25), py + Inches(1.05), Inches(3.25), Inches(0.5))
        s_tf = s_b.text_frame
        p_s = s_tf.paragraphs[0]
        p_s.text = sub
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(18)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE

        # Bullets
        b_b = s4.shapes.add_textbox(px + Inches(0.25), py + Inches(1.65), Inches(3.25), Inches(2.8))
        b_tf = b_b.text_frame
        b_tf.word_wrap = True
        for b_text in bullets:
            p_bullet = b_tf.add_paragraph()
            p_bullet.text = f"•  {b_text}"
            p_bullet.font.name = FONT_MAIN
            p_bullet.font.size = Pt(11.5)
            p_bullet.font.color.rgb = TEXT_MUTED
            p_bullet.space_after = Pt(8)

    add_footer(s4, 4)

    s4.notes_slide.notes_text_frame.text = (
        "SLIDE 4 — ONE PLATFORM. THREE EXPERIENCES.\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Most platforms make the mistake of treating all users like children, or assuming adults want cartoon interfaces. "
        "HumSaathi AI recognizes that needs radically shift across life stages:\n"
        "For Children, it's about sensory predictability and exploratory basics.\n"
        "For Teens, it's about navigating complex peer dynamics and social confidence.\n"
        "For Adults, it's about independent living, workplace meetings, job interviews, and autonomy.\n"
        "Let's look at each portal.\""
    )

    # ==========================================================================
    # SLIDE 05: CHILD PORTAL
    # ==========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_bg(s5)
    add_header(s5, "Early Foundations", "Child Portal: Learn Through Exploration",
               "Designed with occupational therapists in mind: low sensory load, clear contrast, and predictable routines.")

    # Left: Feature Highlights
    add_card(s5, Inches(0.8), Inches(2.1), Inches(4.5), Inches(4.7), bg_color=CARD_DARK, border_color=BORDER_COLOR)
    ch_box = s5.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(3.9), Inches(4.1))
    ch_tf = ch_box.text_frame
    ch_tf.word_wrap = True

    p = ch_tf.paragraphs[0]
    p.text = "CORE PILLARS"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    pillars = [
        ("Simple & Visual", "High-contrast vector assets without sensory-overloading flashes or sudden sounds."),
        ("Predictable Flow", "Consistent three-step activity sequences: Introduce, Practice, Celebrate."),
        ("Bilingual Native Audio", "Voice prompts available in English and Urdu to support home language immersion."),
        ("Zero Pressure", "No countdown timers or punishing fail states. Unlimited gentle retries.")
    ]
    for pil_t, pil_d in pillars:
        p1 = ch_tf.add_paragraph()
        p1.text = f"{pil_t}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(8)

        p2 = ch_tf.add_paragraph()
        p2.text = pil_d
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # Right: Implemented Activities Grid
    act_grid = [
        ("🔤 Letters & Phonics", "English & Urdu alphabet tracing with audio pronunciation"),
        ("🔢 Numbers & Counting", "Predictable counting sequences with visual object arrays"),
        ("🎨 Colors & Shapes", "Geometric recognition and real-world color association"),
        ("🐾 Animals & Nature", "Interactive sounds, visual flashcards, and animal names"),
        ("😊 Emotion Identification", "Facial expression reading with calm social context"),
        ("⏰ Daily Routines", "Step-by-step visual schedules (Brushing, Eating, Bedtime)"),
    ]

    for idx, (title, desc) in enumerate(act_grid):
        rx = Inches(5.6) + (idx % 2) * Inches(3.45)
        ry = Inches(2.1) + (idx // 2) * Inches(1.52)
        add_card(s5, rx, ry, Inches(3.3), Inches(1.38), bg_color=CARD_DARK, border_color=BORDER_COLOR)

        b_box = s5.shapes.add_textbox(rx + Inches(0.15), ry + Inches(0.15), Inches(3.0), Inches(1.1))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        p1 = b_tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

        p2 = b_tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(3)

    add_footer(s5, 5)

    s5.notes_slide.notes_text_frame.text = (
        "SLIDE 5 — CHILD PORTAL (Learn Through Exploration)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"In the Child Portal, we focus on sensory safety and predictable structure. "
        "We've implemented 7 core activity categories: Letters, Numbers, Colors, Shapes, Animals, Emotions, and Daily Routines. "
        "Notice the design philosophy: no countdown timers, no punitive buzzer sounds, and no sudden animations. "
        "Everything is calm, predictable, and supportive of bilingual English and Urdu development.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Mention that activities are fully implemented on both frontend and backend."
    )

    # ==========================================================================
    # SLIDE 06: TEEN + ADULT
    # ==========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_bg(s6)
    add_header(s6, "Maturity Progression", "Teen + Adult: Age-Appropriate Real-World Practice",
               "Moving beyond classroom exercises to real-world autonomy, professional conversations, and social mastery.")

    # Side-by-side Cards
    # TEEN
    add_card(s6, Inches(0.8), Inches(2.1), Inches(5.7), Inches(4.7), bg_color=CARD_DARK, border_color=ACCENT_INDIGO)
    t_badge = add_card(s6, Inches(1.1), Inches(2.35), Inches(2.2), Inches(0.45), bg_color=CARD_ALT, border_color=ACCENT_INDIGO)
    t_box = s6.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(2.2), Inches(0.35))
    p = t_box.text_frame.paragraphs[0]
    p.text = "🧑 TEEN EXPERIENCE"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO

    tb_box = s6.shapes.add_textbox(Inches(1.1), Inches(3.0), Inches(5.1), Inches(3.6))
    tb_tf = tb_box.text_frame
    tb_tf.word_wrap = True

    teen_features = [
        ("Teen Daily Communication Challenge", "Rotating daily micro-tasks: Ask an open-ended question, Express an opinion politely, or Practice saying 'I need help'."),
        ("Peer Conflict & Group Projects", "Interactive roleplay handling disagreements, shared project credit, and making new friends in clubs."),
        ("Context Clues & Reading Passages", "Nuance comprehension: detecting implied meaning and conversational subtleties."),
        ("Authentic Contiguous Streaks", "Real-time streak counting directly connected to database activity history with zero fake stats.")
    ]
    for tf_t, tf_d in teen_features:
        p1 = tb_tf.add_paragraph()
        p1.text = f"🎯  {tf_t}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(4)

        p2 = tb_tf.add_paragraph()
        p2.text = tf_d
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED

    # ADULT
    add_card(s6, Inches(6.8), Inches(2.1), Inches(5.7), Inches(4.7), bg_color=CARD_DARK, border_color=ACCENT_CYAN)
    a_badge = add_card(s6, Inches(7.1), Inches(2.35), Inches(2.2), Inches(0.45), bg_color=CARD_ALT, border_color=ACCENT_CYAN)
    a_box = s6.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(2.2), Inches(0.35))
    p = a_box.text_frame.paragraphs[0]
    p.text = "👨 ADULT EXPERIENCE"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    ab_box = s6.shapes.add_textbox(Inches(7.1), Inches(3.0), Inches(5.1), Inches(3.6))
    ab_tf = ab_box.text_frame
    ab_tf.word_wrap = True

    adult_features = [
        ("Workplace Career Scenarios", "Simulated job interviews, supervisor check-ins, asking for task clarification, and proposing shift swaps."),
        ("Workplace Goals Progress Meter", "Real-time tracking: 'Complete 2 workplace scenarios' with dynamic skill leveling."),
        ("Everyday Autonomy Scenarios", "High-stress daily situations: Bank debit card inquiries, doctor appointments, and restaurant dining."),
        ("Practical Problem Solving", "Contingency management: 'Your bus was cancelled' and 'You received an ambiguous workplace email'.")
    ]
    for af_t, af_d in adult_features:
        p1 = ab_tf.add_paragraph()
        p1.text = f"💼  {af_t}"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(4)

        p2 = ab_tf.add_paragraph()
        p2.text = af_d
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED

    add_footer(s6, 6)

    s6.notes_slide.notes_text_frame.text = (
        "SLIDE 6 — TEEN + ADULT PORTALS (Age-Appropriate Real-World Practice)\n\n"
        "SPEAKING TIME: 35 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Here is where HumSaathi AI truly shines. "
        "In the Teen Portal, we tackle social nuance. Teens receive daily challenges like politely expressing a disagreement or "
        "asking for help when stuck, alongside real streak metrics.\n"
        "In the Adult Portal, the UI shifts to a sleek, professional aesthetic focused on career readiness and independence. "
        "Adults practice job interviews, speaking to managers, handling banking queries, and practical problem-solving like a cancelled bus commute. "
        "This is not gamified fluff; it is real life training.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Highlight that 27 realistic scenarios are built into the platform."
    )

    # ==========================================================================
    # SLIDE 07: AI COMMUNICATION COACH
    # ==========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_bg(s7)
    add_header(s7, "Conversation Engine", "AI That Doesn't Just Answer — It Practices With You",
               "HumSaathi AI engages in realistic turn-by-turn roleplay, evaluating communication and offering immediate coaching.")

    # Left: Turn-by-Turn Simulated Chat Mockup
    add_card(s7, Inches(0.8), Inches(2.1), Inches(6.0), Inches(4.7), bg_color=CARD_DARK, border_color=BORDER_COLOR)

    # Chat Header
    add_card(s7, Inches(1.0), Inches(2.3), Inches(5.6), Inches(0.55), bg_color=CARD_ALT, border_color=BORDER_COLOR)
    ch_h = s7.shapes.add_textbox(Inches(1.2), Inches(2.35), Inches(5.2), Inches(0.45))
    ch_h_tf = ch_h.text_frame
    p = ch_h_tf.paragraphs[0]
    p.text = "🎯 Scenario: Job Interview Practice  |  Role: Hiring Manager"
    p.font.name = FONT_HEADING
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    # Dialogue Bubbles
    dialogue = [
        ("AI Coach", "Welcome! Tell me about a time you handled a difficult deadline.", CARD_ALT, TEXT_WHITE, Inches(2.95), Inches(1.1)),
        ("User (Learner)", "I was nervous, but I asked my lead for priority guidance.", RGBColor(16, 185, 129), TEXT_WHITE, Inches(3.7), Inches(3.6)),
        ("AI Coach", "That shows initiative! How did you update your team after that?", CARD_ALT, TEXT_WHITE, Inches(4.45), Inches(1.1)),
    ]
    for speaker, msg, bg_c, txt_c, top_y, left_x in dialogue:
        add_card(s7, left_x, top_y, Inches(3.8), Inches(0.65), bg_color=bg_c, border_color=BORDER_COLOR)
        d_box = s7.shapes.add_textbox(left_x + Inches(0.15), top_y + Inches(0.08), Inches(3.5), Inches(0.5))
        d_tf = d_box.text_frame
        d_tf.word_wrap = True
        p1 = d_tf.paragraphs[0]
        p1.text = speaker.upper()
        p1.font.size = Pt(8.5)
        p1.font.bold = True
        p1.font.color.rgb = ACCENT_CYAN if speaker == "AI Coach" else RGBColor(209, 250, 229)
        p2 = d_tf.add_paragraph()
        p2.text = msg
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = txt_c

    # Right: Structured Post-Practice Evaluation Card
    add_card(s7, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.7), bg_color=CARD_DARK, border_color=ACCENT_CYAN)

    ev_h = s7.shapes.add_textbox(Inches(7.4), Inches(2.3), Inches(4.8), Inches(0.45))
    p = ev_h.text_frame.paragraphs[0]
    p.text = "STRUCTURED EVALUATION REPORT"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    # 3 Metrics
    metrics = [
        ("Clarity Score", "85%", ACCENT_GREEN, Inches(2.8)),
        ("Relevance Score", "90%", ACCENT_CYAN, Inches(3.45)),
        ("Confidence Score", "80%", ACCENT_INDIGO, Inches(4.1)),
    ]
    for m_label, m_val, m_col, m_top in metrics:
        add_card(s7, Inches(7.4), m_top, Inches(4.8), Inches(0.55), bg_color=CARD_ALT, border_color=BORDER_COLOR)
        m_b = s7.shapes.add_textbox(Inches(7.6), m_top + Inches(0.1), Inches(4.4), Inches(0.35))
        m_tf = m_b.text_frame
        p = m_tf.paragraphs[0]
        p.text = f"{m_label}:  {m_val}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = m_col

    # Feedback callout
    add_card(s7, Inches(7.4), Inches(4.85), Inches(4.8), Inches(1.75), bg_color=CARD_ALT, border_color=ACCENT_CYAN)
    fb_b = s7.shapes.add_textbox(Inches(7.55), Inches(4.95), Inches(4.5), Inches(1.55))
    fb_tf = fb_b.text_frame
    fb_tf.word_wrap = True

    p = fb_tf.paragraphs[0]
    p.text = "💡 SUGGESTED ALTERNATIVE RESPONSE"
    p.font.name = FONT_MAIN
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p_alt = fb_tf.add_paragraph()
    p_alt.text = "\"When deadlines tightened, I immediately consulted my lead to prioritize high-impact deliverables, keeping everyone synchronized.\""
    p_alt.font.name = FONT_MAIN
    p_alt.font.size = Pt(10.5)
    p_alt.font.italic = True
    p_alt.font.color.rgb = TEXT_WHITE
    p_alt.space_before = Pt(4)

    add_footer(s7, 7)

    s7.notes_slide.notes_text_frame.text = (
        "SLIDE 7 — AI COMMUNICATION COACH (Practices With You)\n\n"
        "SPEAKING TIME: 40 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"This is our standout technical feature. HumSaathi AI does not just answer questions like a generic chatbot. "
        "It steps into character — as an interviewer, a coworker, or a peer — and conducts realistic turn-by-turn roleplay.\n"
        "Look at what happens after the conversation:\n"
        "The AI produces a multi-dimensional diagnostic evaluating Clarity, Relevance, and Confidence. "
        "Crucially, it provides a 'Suggested Alternative Response', showing learners exactly how to rephrase their thoughts "
        "more naturally and professionally next time.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Emphasize that it teaches communication *mechanics*, not just facts."
    )

    # ==========================================================================
    # SLIDE 08: PRODUCTION AI ARCHITECTURE
    # ==========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_bg(s8)
    add_header(s8, "Deep-Dive Engineering", "Production AI Architecture",
               "A multi-tier architecture combining referent anchor tracking, persona grounding, and structured evaluation.")

    arch_layers = [
        ("1. CLIENT LAYER", "React 18 + Vite SPA", "Web Speech API audio stream, trilingual I18n UI, sensory mode toggle", ACCENT_GREEN, Inches(0.8)),
        ("2. GATEWAY LAYER", "FastAPI Python 3.11", "ASGI serverless execution, JWT token validation, rate-limiting", ACCENT_CYAN, Inches(2.78)),
        ("3. CONTEXT & PERSONA", "Grounding Engine", "Referent anchor tracking, intent classification, multi-turn history", ACCENT_INDIGO, Inches(4.76)),
        ("4. LLM & FALLBACK", "Gemini 2.0 / OpenAI", "Temperature-calibrated persona prompting with deterministic rule fallback", ACCENT_PURPLE, Inches(6.74)),
        ("5. EVALUATION", "Scoring Engine", "Clarity, relevance, confidence metrics & alternative phrasing generation", ACCENT_ORANGE, Inches(8.72)),
        ("6. PERSISTENCE", "Database & Analytics", "Persona-isolated progress, contiguous streaks, weekly parent insights", ACCENT_GREEN, Inches(10.7)),
    ]

    for title, sub, desc, col, lx in arch_layers:
        add_card(s8, lx, Inches(2.2), Inches(1.85), Inches(4.0), bg_color=CARD_DARK, border_color=col)
        
        # Pill
        add_card(s8, lx + Inches(0.1), Inches(2.35), Inches(1.65), Inches(0.35), bg_color=CARD_ALT, border_color=col)
        p_b = s8.shapes.add_textbox(lx + Inches(0.1), Inches(2.38), Inches(1.65), Inches(0.3))
        p_tf = p_b.text_frame
        p = p_tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(8.5)
        p.font.bold = True
        p.font.color.rgb = col

        # Sub
        s_b = s8.shapes.add_textbox(lx + Inches(0.1), Inches(2.8), Inches(1.65), Inches(0.6))
        s_tf = s_b.text_frame
        s_tf.word_wrap = True
        p_s = s_tf.paragraphs[0]
        p_s.text = sub
        p_s.font.name = FONT_HEADING
        p_s.font.size = Pt(12.5)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE

        # Desc
        d_b = s8.shapes.add_textbox(lx + Inches(0.1), Inches(3.45), Inches(1.65), Inches(2.6))
        d_tf = d_b.text_frame
        d_tf.word_wrap = True
        p_d = d_tf.paragraphs[0]
        p_d.text = desc
        p_d.font.name = FONT_MAIN
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_MUTED

    # Bottom Technical Note
    bot_card = add_card(s8, Inches(0.8), Inches(6.35), Inches(11.733), Inches(0.55), bg_color=CARD_ALT, border_color=BORDER_COLOR)
    bot_b = s8.shapes.add_textbox(Inches(1.0), Inches(6.42), Inches(11.333), Inches(0.4))
    p = bot_b.text_frame.paragraphs[0]
    p.text = "⚡ PRODUCTION RESILIENCE: If external AI APIs experience latency or outages, our rule-based policy engine ensures 100% uninterrupted offline coaching."
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    add_footer(s8, 8)

    s8.notes_slide.notes_text_frame.text = (
        "SLIDE 8 — PRODUCTION AI ARCHITECTURE (How The AI Works)\n\n"
        "SPEAKING TIME: 40 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"For our technical judges: how does this actually work under the hood?\n"
        "Our frontend in React 18 streams user voice and text into a Python FastAPI backend.\n"
        "Before sending any request to an LLM, it passes through our Context and Persona Grounding Engine in context_builder.py. "
        "This resolves referent pronouns, classifies intent, attaches the user's communication profile, and grounds the conversation.\n"
        "We support Google Gemini 2.0 and OpenAI models. But notice our resiliency layer: if API limits or network drops occur, "
        "our deterministic conversation engine steps in seamlessly, ensuring zero downtime for the learner.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Demonstrates production-grade design and thoughtful offline fallback."
    )

    # ==========================================================================
    # SLIDE 09: ADAPTIVE LEARNING LOOP
    # ==========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_bg(s9)
    add_header(s9, "Continuous Calibration", "Adaptive Learning Loop",
               "Every interaction makes the next one smarter. HumSaathi AI continuously calibrates challenge and support.")

    loop_steps = [
        ("1. ASSESS", "Baseline Diagnostics", "Learner completes lightweight baseline evaluation to identify strengths & sensory preferences.", ACCENT_GREEN),
        ("2. PRACTICE", "Interactive Execution", "Learner engages in voice/text communication scenarios and structured skill modules.", ACCENT_CYAN),
        ("3. OBSERVE", "Behavioral Signals", "System observes response latency, sentence complexity, and vocabulary retention in real time.", ACCENT_INDIGO),
        ("4. EVALUATE", "Tri-Metric Scoring", "Multi-dimensional scoring computes Clarity, Relevance, and Confidence scores instantly.", ACCENT_PURPLE),
        ("5. ADAPT", "Cognitive Adjustment", "Engine recalibrates scenario difficulty, sensory stimulation levels, and pacing thresholds.", ACCENT_ORANGE),
        ("6. RECOMMEND", "Next Best Practice", "Generates the 'Next Best Practice' card directly on the learner's dashboard.", ACCENT_GREEN),
    ]

    for i, (title, sub, body, col) in enumerate(loop_steps):
        col_idx = i % 3
        row_idx = i // 3
        cx = Inches(0.8) + col_idx * Inches(3.99)
        cy = Inches(2.1) + row_idx * Inches(2.4)

        add_card(s9, cx, cy, Inches(3.75), Inches(2.2), bg_color=CARD_DARK, border_color=col)

        # Header Pill
        add_card(s9, cx + Inches(0.2), cy + Inches(0.2), Inches(3.35), Inches(0.4), bg_color=CARD_ALT, border_color=col)
        h_b = s9.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.24), Inches(3.35), Inches(0.35))
        p = h_b.text_frame.paragraphs[0]
        p.text = f"{title}  ·  {sub}"
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        # Body
        b_b = s9.shapes.add_textbox(cx + Inches(0.2), cy + Inches(0.7), Inches(3.35), Inches(1.4))
        b_tf = b_b.text_frame
        b_tf.word_wrap = True
        p_b = b_tf.paragraphs[0]
        p_b.text = body
        p_b.font.name = FONT_MAIN
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = TEXT_MUTED

    add_footer(s9, 9)

    s9.notes_slide.notes_text_frame.text = (
        "SLIDE 9 — ADAPTIVE LEARNING LOOP\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"HumSaathi AI is an adaptive loop: Assess, Practice, Observe, Evaluate, Adapt, and Recommend.\n"
        "When a learner struggles with expressing disagreement in a teen club scenario, the engine notices the hesitation. "
        "Instead of penalizing them, it adapts: it suggests an intermediate roleplay on active listening, and updates the "
        "dashboard's 'Next Best Practice' banner.\n"
        "The system grows with the learner day by day.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Explains the recommendation algorithm and personalized progression."
    )

    # ==========================================================================
    # SLIDE 10: MULTIMODAL COMMUNICATION
    # ==========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_bg(s10)
    add_header(s10, "Inclusive Interfaces", "Multimodal Communication: Text + Voice + AI",
               "Supporting continuous voice interaction, speech recognition, and native trilingual localization.")

    # Left: Multimodal Architecture Flow
    add_card(s10, Inches(0.8), Inches(2.1), Inches(6.8), Inches(4.7), bg_color=CARD_DARK, border_color=BORDER_COLOR)

    flow_items = [
        ("🎙️ Voice & Speech Recognition", "Web Speech API captures spoken audio, providing accessible real-world phone call simulation."),
        ("💬 Real-Time Text Input", "Comfortable, low-anxiety text interface for learners who prefer typing or AAC assistive tools."),
        ("🧠 Trilingual AI Understanding", "Context engine auto-detects and processes linguistic nuances across 3 languages."),
        ("🔊 Articulate Speech Synthesis", "Calm, natural text-to-speech output delivering empathetic conversational pacing."),
    ]
    fb_box = s10.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(6.2), Inches(4.3))
    fb_tf = fb_box.text_frame
    fb_tf.word_wrap = True
    p = fb_tf.paragraphs[0]
    p.text = "MULTIMODAL INTERACTION PIPELINE"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    for title, desc in flow_items:
        p1 = fb_tf.add_paragraph()
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(6)

        p2 = fb_tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # Right: Language Cards
    langs = [
        ("ENGLISH", "en", "Standard international phrasing and professional workplace terminology", ACCENT_CYAN),
        ("URDU (اردو)", "ur", "Full native Right-to-Left (RTL) typography with cultural phrasing", ACCENT_GREEN),
        ("ROMAN URDU", "ur_rm", "Phonetic conversational vernacular widely used across South Asia", ACCENT_INDIGO),
    ]
    for i, (name, code, desc, col) in enumerate(langs):
        ly = Inches(2.1) + i * Inches(1.58)
        add_card(s10, Inches(7.9), ly, Inches(4.6), Inches(1.45), bg_color=CARD_DARK, border_color=col)
        
        l_box = s10.shapes.add_textbox(Inches(8.1), ly + Inches(0.15), Inches(4.2), Inches(1.15))
        l_tf = l_box.text_frame
        l_tf.word_wrap = True

        p1 = l_tf.paragraphs[0]
        p1.text = f"{name}  [{code}]"
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = l_tf.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(3)

    add_footer(s10, 10)

    s10.notes_slide.notes_text_frame.text = (
        "SLIDE 10 — MULTIMODAL COMMUNICATION (Text + Voice + AI)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Real communication happens through both speech and text. HumSaathi AI integrates the browser's Web Speech API "
        "to deliver simulated telephone and in-person voice conversations.\n"
        "Furthermore, we recognize regional accessibility: HumSaathi is natively trilingual. It supports English, native Urdu "
        "with complete RTL typography, and Roman Urdu — ensuring learners across South Asia and global diaspora communities "
        "can practice in the dialect they speak at home.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Highlight cultural inclusivity and real-world applicability."
    )

    # ==========================================================================
    # SLIDE 11: PARENT COMPANION
    # ==========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_bg(s11)
    add_header(s11, "Family Ecosystem", "Parent Companion: Understand Progress — Not Control It",
               "A privacy-first portal giving caregivers actionable visibility into communication growth without invasive micromanagement.")

    # 4 Dimension Cards
    dims = [
        ("👋 GREETING", "Social Warmth", "Measures ability to initiate and acknowledge conversational openings gracefully.", ACCENT_GREEN),
        ("💬 ANSWERING QUESTIONS", "Response Accuracy", "Evaluates clarity, directness, and relevance when asked for information.", ACCENT_CYAN),
        ("❓ ASKING QUESTIONS", "Curiosity & Engagement", "Tracks proactive inquiry, requesting clarification, and conversational balance.", ACCENT_INDIGO),
        ("🌟 CONFIDENCE", "Emotional Stability", "Measures smooth conversational pacing, reduced hesitation, and natural flow.", ACCENT_PURPLE),
    ]
    for i, (title, sub, desc, col) in enumerate(dims):
        dx = Inches(0.8) + i * Inches(2.98)
        add_card(s11, dx, Inches(2.1), Inches(2.78), Inches(2.6), bg_color=CARD_DARK, border_color=col)

        b_box = s11.shapes.add_textbox(dx + Inches(0.15), Inches(2.25), Inches(2.48), Inches(2.3))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        p1 = b_tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = b_tf.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(2)

        p3 = b_tf.add_paragraph()
        p3.text = desc
        p3.font.name = FONT_MAIN
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(6)

    # Bottom Split: Weekly AI Report & Privacy Security
    add_card(s11, Inches(0.8), Inches(4.9), Inches(6.8), Inches(1.9), bg_color=CARD_DARK, border_color=BORDER_COLOR)
    w_box = s11.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(6.4), Inches(1.7))
    w_tf = w_box.text_frame
    w_tf.word_wrap = True
    p = w_tf.paragraphs[0]
    p.text = "📊 AI WEEKLY OBSERVATION DIGEST"
    p.font.name = FONT_MAIN
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p_w = w_tf.add_paragraph()
    p_w.text = "• Highlights verified developmental wins: 'Demonstrated 25% faster clarification turnaround.'\n• Suggests real-world family activities: 'Ask Zayd to order at dinner tonight to reinforce practice.'\n• Never displays raw transcripts, preserving learner dignity and autonomy."
    p_w.font.name = FONT_MAIN
    p_w.font.size = Pt(11)
    p_w.font.color.rgb = TEXT_MUTED
    p_w.space_before = Pt(4)

    add_card(s11, Inches(7.9), Inches(4.9), Inches(4.6), Inches(1.9), bg_color=CARD_DARK, border_color=BORDER_COLOR)
    p_box = s11.shapes.add_textbox(Inches(8.1), Inches(5.0), Inches(4.2), Inches(1.7))
    p_tf = p_box.text_frame
    p_tf.word_wrap = True
    p = p_tf.paragraphs[0]
    p.text = "🔒 CAREGIVER PRIVACY & ACCESS"
    p.font.name = FONT_MAIN
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p_sec = p_tf.add_paragraph()
    p_sec.text = "• Protected by a localized 4-Digit Parent PIN gate.\n• Complete persona data isolation prevents cross-profile leaks.\n• Meets international standards for child privacy and accessibility."
    p_sec.font.name = FONT_MAIN
    p_sec.font.size = Pt(11)
    p_sec.font.color.rgb = TEXT_MUTED
    p_sec.space_before = Pt(4)

    add_footer(s11, 11)

    s11.notes_slide.notes_text_frame.text = (
        "SLIDE 11 — PARENT COMPANION (Understand Progress — Not Control It)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Parents of neurodiverse children often feel blind: 'Is my child improving? How do I help them at home?'\n"
        "The Parent Portal solves this with four core communication metrics: Greeting, Answering, Asking, and Confidence.\n"
        "Crucially, we maintain privacy: parents do not read private transcripts. Instead, our AI Companion synthesizes "
        "actionable insights, telling parents: 'Your child mastered asking for help this week. Encourage them to order at the restaurant tonight.' "
        "It turns technology into real family connection.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Emphasize privacy and constructive parental empowerment."
    )

    # ==========================================================================
    # SLIDE 12: TECHNOLOGY STACK
    # ==========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_bg(s12)
    add_header(s12, "Production Engineering", "Full-Stack AI Ecosystem",
               "Built on a modern, decoupled architecture designed for high throughput, serverless efficiency, and rapid scale.")

    tech_blocks = [
        ("FRONTEND", "React 18 & Vite", "SPA architecture, React Router v7, zero heavy UI frameworks, native CSS variables, responsive accessible layout.", ACCENT_CYAN),
        ("BACKEND API", "FastAPI & Python 3.11", "Asynchronous ASGI serverless endpoints, Pydantic v2 data validation, Uvicorn high-concurrency runtime.", ACCENT_GREEN),
        ("AI BRAIN", "Gemini 2.0 & OpenAI", "Dynamic prompt chaining, intent classification, persona policy enforcement, and deterministic rule fallbacks.", ACCENT_PURPLE),
        ("PERSISTENCE", "SQLAlchemy & SQLite/Postgres", "Robust relational schema supporting session logs, multidimensional evaluations, and real-time streak computation.", ACCENT_INDIGO),
        ("SPEECH ENGINE", "Web Speech API", "Continuous client-side speech recognition and synthesis with zero external third-party speech vendor latency.", ACCENT_ORANGE),
        ("INFRASTRUCTURE", "Vercel Serverless & CDN", "Continuous Git integration, serverless function rewrites via vercel.json, 100% production uptime.", ACCENT_CYAN),
    ]

    for i, (cat, tech, desc, col) in enumerate(tech_blocks):
        col_idx = i % 3
        row_idx = i // 3
        tx = Inches(0.8) + col_idx * Inches(3.99)
        ty = Inches(2.1) + row_idx * Inches(2.4)

        add_card(s12, tx, ty, Inches(3.75), Inches(2.2), bg_color=CARD_DARK, border_color=BORDER_COLOR)

        b_box = s12.shapes.add_textbox(tx + Inches(0.2), ty + Inches(0.18), Inches(3.35), Inches(1.8))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        p1 = b_tf.paragraphs[0]
        p1.text = cat
        p1.font.name = FONT_MAIN
        p1.font.size = Pt(9.5)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = b_tf.add_paragraph()
        p2.text = tech
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(2)

        p3 = b_tf.add_paragraph()
        p3.text = desc
        p3.font.name = FONT_MAIN
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(6)

    add_footer(s12, 12)

    s12.notes_slide.notes_text_frame.text = (
        "SLIDE 12 — TECHNOLOGY STACK (Full-Stack AI Ecosystem)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Here is our actual technology stack:\n"
        "We chose React 18 and Vite for near-instant rendering and low sensory load.\n"
        "Our backend runs Python 3.11 with FastAPI, providing lightning-fast asynchronous ASGI response times.\n"
        "For AI, we integrate Google Gemini 2.0 with strict temperature calibration and fallback mechanisms.\n"
        "Everything is deployed live on Vercel's global edge network at hum-saathi-ai.vercel.app with automated CI/CD.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Emphasize that this is already live in production, not a mock repository."
    )

    # ==========================================================================
    # SLIDE 13: SECURITY & ACCESS CONTROL
    # ==========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_bg(s13)
    add_header(s13, "Data Protection", "Designed With Privacy in Mind",
               "Neurodiverse learners deserve the highest standard of data protection, role isolation, and ethical AI safety.")

    sec_cards = [
        ("🔐 SECURE AUTHENTICATION", "Dual Authentication Flow",
         "Robust JWT token-based authentication paired with Google OAuth 2.0 Identity Services for secure, passwordless login.", ACCENT_CYAN),
        ("🛡️ PERSONA DATA ISOLATION", "Zero Cross-Contamination",
         "Strict backend database filtering guarantees that teen or adult scenarios and progress metrics never leak into child accounts.", ACCENT_GREEN),
        ("🔑 CAREGIVER PIN GATE", "Localized Security",
         "Parent Portal access is gated behind a local encrypted 4-digit PIN, ensuring children cannot accidentally modify profiles.", ACCENT_INDIGO),
        ("👮 ADMIN ROLE CONTROL", "Comprehensive RBAC",
         "Enterprise Role-Based Access Control allows verified administrators to monitor system health, audit logs, and AI behavior.", ACCENT_PURPLE),
    ]

    for i, (title, sub, desc, col) in enumerate(sec_cards):
        sx = Inches(0.8) + i * Inches(2.98)
        add_card(s13, sx, Inches(2.1), Inches(2.78), Inches(4.7), bg_color=CARD_DARK, border_color=col)

        b_box = s13.shapes.add_textbox(sx + Inches(0.2), Inches(2.35), Inches(2.38), Inches(4.2))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        p1 = b_tf.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_MAIN
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = col

        p2 = b_tf.add_paragraph()
        p2.text = sub
        p2.font.name = FONT_HEADING
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        p2.space_before = Pt(4)

        p3 = b_tf.add_paragraph()
        p3.text = desc
        p3.font.name = FONT_MAIN
        p3.font.size = Pt(11)
        p3.font.color.rgb = TEXT_MUTED
        p3.space_before = Pt(8)

    add_footer(s13, 13)

    s13.notes_slide.notes_text_frame.text = (
        "SLIDE 13 — SECURITY & ACCESS CONTROL (Privacy in Mind)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"When dealing with vulnerable populations, safety and privacy cannot be an afterthought.\n"
        "HumSaathi AI implements strict persona data isolation — an adult's workplace metrics and a child's foundational phonics "
        "are physically segregated at the database query level.\n"
        "We support Google OAuth 2.0 and JWT auth, protect the Parent Portal with an encrypted PIN, and provide comprehensive "
        "Role-Based Access Control with audit logs in the Admin Center.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Conveys enterprise readiness and ethical AI compliance."
    )

    # ==========================================================================
    # SLIDE 14: WHY HUMSAATHI AI?
    # ==========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_bg(s14)
    add_header(s14, "Competitive Differentiation", "Why HumSaathi AI? More Than an AI Chatbot",
               "How HumSaathi AI fundamentally differs from standard generative AI chatbots and traditional e-learning apps.")

    # Table Card
    add_card(s14, Inches(0.8), Inches(2.1), Inches(11.733), Inches(4.7), bg_color=CARD_DARK, border_color=BORDER_COLOR)

    table_shape = s14.shapes.add_table(7, 3, Inches(1.1), Inches(2.3), Inches(11.133), Inches(4.3))
    table = table_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(4.333)

    headers = ["DIMENSION", "GENERIC AI CHATBOTS (e.g. ChatGPT)", "HUMSAATHI AI"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_ALT
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = FONT_HEADING
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN if j == 2 else TEXT_WHITE

    rows = [
        ("Learning Structure", "Unbounded, open-ended prompt box (high cognitive load)", "27 goal-oriented, structured, scaffolded scenarios"),
        ("Lifespan Personalization", "Single generic tone for all queries", "Differentiated Child, Teen, and Adult portal workflows"),
        ("Sensory Ergonomics", "Monochrome, text-heavy, distracting UI", "Calm, low-sensory palette, high contrast, predictable flow"),
        ("Communication Evaluation", "Vague 'Good job' or generic text summaries", "Objective metrics: Clarity, Relevance, and Confidence"),
        ("Actionable Feedback", "No explicit guidance on alternative phrasing", "Offers direct 'Suggested Alternative Response' coaching"),
        ("Caregiver Visibility", "Zero oversight or family coaching integration", "Privacy-first Parent Companion with weekly actionable tips"),
    ]

    for i, (dim, gen, hum) in enumerate(rows):
        r_idx = i + 1
        c0 = table.cell(r_idx, 0)
        c0.fill.solid()
        c0.fill.fore_color.rgb = CARD_DARK
        p0 = c0.text_frame.paragraphs[0]
        p0.text = dim
        p0.font.bold = True
        p0.font.size = Pt(10.5)
        p0.font.color.rgb = TEXT_WHITE

        c1 = table.cell(r_idx, 1)
        c1.fill.solid()
        c1.fill.fore_color.rgb = CARD_DARK
        p1 = c1.text_frame.paragraphs[0]
        p1.text = gen
        p1.font.size = Pt(10)
        p1.font.color.rgb = TEXT_MUTED

        c2 = table.cell(r_idx, 2)
        c2.fill.solid()
        c2.fill.fore_color.rgb = CARD_ALT
        p2 = c2.text_frame.paragraphs[0]
        p2.text = f"✓  {hum}"
        p2.font.bold = True
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = ACCENT_GREEN

    add_footer(s14, 14)

    s14.notes_slide.notes_text_frame.text = (
        "SLIDE 14 — WHY HUMSAATHI AI? (Competitive Differentiation)\n\n"
        "SPEAKING TIME: 35 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Judges often ask: 'Why not just use ChatGPT or Duolingo?'\n"
        "Here is the fundamental difference: an open prompt box is overwhelming for someone with neurodiversity or anxiety. "
        "HumSaathi AI provides scaffolded structure.\n"
        "Unlike generic LLMs, HumSaathi evaluates Clarity, Relevance, and Confidence, provides suggested alternative phrasing, "
        "and is ergonomically designed for low sensory load. It is an educational health intervention, not a text toy.\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Strong competitive moat established across ergonomics, evaluation, and caregiver integration."
    )

    # ==========================================================================
    # SLIDE 15: FUTURE + CLOSING
    # ==========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_bg(s15, BG_DARKER)

    # Top Header
    eb_box = s15.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.733), Inches(0.4))
    eb_tf = eb_box.text_frame
    p = eb_tf.paragraphs[0]
    p.text = "THE FUTURE OF PERSONALIZED COMMUNICATION LEARNING"
    p.font.name = FONT_MAIN
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    # Roadmap Grid
    road_items = [
        ("📹 Multimodal Video & Eye Contact", "Future Scope", "Computer vision facial sentiment & non-verbal communication coaching."),
        ("👩‍⚕️ Therapist & Clinic Dashboard", "Future Scope", "Direct telehealth portal for speech pathologists and behavioral therapists."),
        ("📱 Native Mobile Application", "Future Scope", "Offline-first React Native mobile apps for iOS and Android devices."),
        ("🌐 Global Vernacular Expansion", "Future Scope", "Arabic, Spanish, Bengali, and Hindi dialect adaptation modules."),
    ]
    for i, (rtitle, rtag, rdesc) in enumerate(road_items):
        rx = Inches(0.8) + i * Inches(2.98)
        add_card(s15, rx, Inches(1.15), Inches(2.78), Inches(1.65), bg_color=CARD_DARK, border_color=BORDER_COLOR)

        b_box = s15.shapes.add_textbox(rx + Inches(0.15), Inches(1.25), Inches(2.48), Inches(1.45))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        p1 = b_tf.paragraphs[0]
        p1.text = rtitle
        p1.font.name = FONT_HEADING
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

        p2 = b_tf.add_paragraph()
        p2.text = f"[{rtag}]  {rdesc}"
        p2.font.name = FONT_MAIN
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(3)

    # Center Hero Box for Memorable Closing
    add_card(s15, Inches(1.5), Inches(3.1), Inches(10.333), Inches(3.7), bg_color=CARD_DARK, border_color=ACCENT_CYAN)

    if has_logo:
        try:
            s15.shapes.add_picture(logo_path, Inches(5.9), Inches(3.3), width=Inches(1.5))
        except Exception:
            pass

    c_box = s15.shapes.add_textbox(Inches(2.0), Inches(4.8), Inches(9.333), Inches(0.6))
    c_tf = c_box.text_frame
    p_c = c_tf.paragraphs[0]
    p_c.text = "HumSaathi AI  ·  Learn. Communicate. Grow."
    p_c.alignment = PP_ALIGN.CENTER
    p_c.font.name = FONT_HEADING
    p_c.font.size = Pt(22)
    p_c.font.bold = True
    p_c.font.color.rgb = ACCENT_CYAN

    q_box = s15.shapes.add_textbox(Inches(2.0), Inches(5.4), Inches(9.333), Inches(1.1))
    q_tf = q_box.text_frame
    q_tf.word_wrap = True
    p_q = q_tf.paragraphs[0]
    p_q.text = "“HumSaathi AI doesn't just teach users what to learn —\nit helps them learn how to communicate, practice, and grow.”"
    p_q.alignment = PP_ALIGN.CENTER
    p_q.font.name = FONT_MAIN
    p_q.font.size = Pt(17)
    p_q.font.bold = True
    p_q.font.italic = True
    p_q.font.color.rgb = TEXT_WHITE

    add_footer(s15, 15)

    s15.notes_slide.notes_text_frame.text = (
        "SLIDE 15 — FUTURE + CLOSING (Learn. Communicate. Grow.)\n\n"
        "SPEAKING TIME: 30 Seconds\n\n"
        "WHAT TO SAY:\n"
        "\"Looking ahead, our roadmap includes therapist telehealth dashboards, offline-first mobile apps, and non-verbal expression coaching.\n"
        "To conclude:\n"
        "HumSaathi AI doesn't just teach users what to learn — it helps them learn how to communicate, practice, and grow.\n"
        "Thank you. The platform is live right now at hum-saathi-ai.vercel.app, and we welcome your questions!\"\n\n"
        "KEY POINT FOR JUDGES:\n"
        "- Strong, memorable, emotional yet technical conclusion."
    )

    # Save presentation
    output_filename = "HumSaathi_AI_Hackathon_Pitch.pptx"
    prs.save(output_filename)
    print(f"SUCCESS: Generated {output_filename} with 15 premium widescreen slides!")

if __name__ == "__main__":
    create_deck()
